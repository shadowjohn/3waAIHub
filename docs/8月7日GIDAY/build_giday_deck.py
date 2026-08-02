#!/usr/bin/env python3
from __future__ import annotations

import math
import random
import sys
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/var/www/html/3waAIHub/docs/8月7日GIDAY")
ASSETS = ROOT / "assets"
LIGHT_THEME = "--light" in sys.argv
OUT = ROOT / (
    "GIS集團-AI服務平台_GIDAY_2026-08-07_淡色版.pptx"
    if LIGHT_THEME
    else "GIS集團-AI服務平台_GIDAY_2026-08-07.pptx"
)

W, H = 13.333, 7.5

if LIGHT_THEME:
    BG = "F4F8FC"
    BG2 = "E8F1F7"
    PANEL = "FFFFFF"
    LINE = "BDD2E1"
    TEXT = "10243A"
    MUTED = "526B80"
    CYAN = "007FA8"
    BLUE = "2563EB"
    TEAL = "008A72"
    GREEN = "0C8C63"
    AMBER = "B96600"
    RED = "C63D55"
    GRID = "DDEAF2"
    BAR_BG = "D5E4EE"
    RISK_PANEL = "FFF4F5"
    SUCCESS_PANEL = "EFFAF6"
else:
    BG = "07111F"
    BG2 = "0B1828"
    PANEL = "102235"
    LINE = "24435C"
    TEXT = "F5F9FC"
    MUTED = "A8BED2"
    CYAN = "37D7FF"
    BLUE = "2F73FF"
    TEAL = "35E6B1"
    GREEN = "4BE49A"
    AMBER = "FFB84A"
    RED = "FF6277"
    GRID = "0E2638"
    BAR_BG = "153148"
    RISK_PANEL = "171722"
    SUCCESS_PANEL = "0D2028"
WHITE = "FFFFFF"

CJK = "Noto Sans CJK TC"
DISPLAY = "Tektur"
LATIN = "Instrument Sans"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_alpha(shape, alpha: int) -> None:
    """alpha: 0 transparent, 100 opaque."""
    fill = shape.fill
    fill.solid()
    solid = fill._xPr.solidFill
    for child in list(solid):
        for alpha_node in child.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha"
        ):
            child.remove(alpha_node)
        alpha_node = child.makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
            {"val": str(max(0, min(100000, alpha * 1000)))},
        )
        child.append(alpha_node)


def add_rect(
    slide,
    x,
    y,
    w,
    h,
    fill=PANEL,
    line=None,
    radius=True,
    alpha=100,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    set_alpha(shape, alpha)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, x, y, d, fill=PANEL, line=None, alpha=100):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    set_alpha(shape, alpha)
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=20,
    color=TEXT,
    bold=False,
    font=CJK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
    line_spacing=1.05,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = Inches(margin)
    box.text_frame.margin_right = Inches(margin)
    box.text_frame.margin_top = Inches(margin)
    box.text_frame.margin_bottom = Inches(margin)
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_runs(
    slide,
    runs,
    x,
    y,
    w,
    h,
    size=20,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    line_spacing=1.02,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = True
    box.text_frame.vertical_anchor = valign
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for item in runs:
        run = p.add_run()
        run.text = item["text"]
        run.font.name = item.get("font", CJK)
        run.font.size = Pt(item.get("size", size))
        run.font.bold = item.get("bold", False)
        run.font.color.rgb = rgb(item.get("color", TEXT))
    return box


def add_paragraphs(
    slide,
    paragraphs,
    x,
    y,
    w,
    h,
    size=18,
    color=TEXT,
    bullet=False,
    spacing=8,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    box.text_frame.word_wrap = True
    for index, item in enumerate(paragraphs):
        p = box.text_frame.paragraphs[0] if index == 0 else box.text_frame.add_paragraph()
        p.text = item
        p.font.name = CJK
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
        p.space_after = Pt(spacing)
        p.line_spacing = 1.05
        if bullet:
            p.text = "•  " + item
    return box


def add_line(slide, x1, y1, x2, y2, color=CYAN, width=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    return line


def add_chevron(slide, x, y, w=0.32, h=0.55, color=CYAN):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.line.fill.background()
    return shape


def add_image_crop(slide, path: Path, x, y, w, h):
    with Image.open(path) as im:
        iw, ih = im.size
    pic = slide.shapes.add_picture(
        str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h)
    )
    image_ratio = iw / ih
    frame_ratio = w / h
    if image_ratio > frame_ratio:
        crop = 1 - frame_ratio / image_ratio
        pic.crop_left = crop / 2
        pic.crop_right = crop / 2
    else:
        crop = 1 - image_ratio / frame_ratio
        pic.crop_top = crop / 2
        pic.crop_bottom = crop / 2
    return pic


def add_header(slide, page, title, section):
    add_text(
        slide,
        f"{page:02d}  {section.upper()}",
        0.68,
        0.25,
        3.0,
        0.28,
        10,
        CYAN,
        bold=True,
        font=DISPLAY,
    )
    add_text(slide, title, 0.68, 0.57, 12.0, 0.65, 29, TEXT, bold=True)
    add_line(slide, 0.68, 1.28, 12.65, 1.28, LINE, 1)


def add_footer(slide, page, note=""):
    if note:
        add_text(slide, note, 0.68, 7.11, 10.7, 0.2, 8, MUTED)
    add_text(
        slide,
        f"{page:02d}",
        12.18,
        7.05,
        0.45,
        0.25,
        9,
        CYAN,
        bold=True,
        font=DISPLAY,
        align=PP_ALIGN.RIGHT,
    )


def add_bg(slide, color=BG):
    add_rect(slide, 0, 0, W, H, fill=color, radius=False)
    for x in (0.65, 4.65, 8.65, 12.65):
        add_line(slide, x, 0, x, H, GRID, 0.5)
    add_line(slide, 0, 6.98, W, 6.98, GRID, 0.5)


def add_metric(slide, x, y, w, value, label, color=CYAN, sub=None):
    add_text(slide, value, x, y, w, 0.68, 32, color, bold=True, font=DISPLAY)
    add_text(slide, label, x, y + 0.72, w, 0.32, 14, TEXT, bold=True)
    if sub:
        add_text(slide, sub, x, y + 1.02, w, 0.35, 10, MUTED)


def add_node(slide, x, y, w, h, title, lines, accent=CYAN, monogram="N"):
    add_rect(slide, x, y, w, h, fill=PANEL, line=LINE)
    add_circle(slide, x + 0.18, y + 0.18, 0.54, fill=BG2, line=accent)
    add_text(
        slide,
        monogram,
        x + 0.18,
        y + 0.18,
        0.54,
        0.54,
        16,
        accent,
        bold=True,
        font=DISPLAY,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, title, x + 0.86, y + 0.18, w - 1.04, 0.42, 13, TEXT, bold=True)
    add_text(slide, "\n".join(lines), x + 0.86, y + 0.64, w - 1.04, h - 0.76, 11, MUTED)


def add_person(slide, x, y, scale=1.0, accent=CYAN, fill=PANEL):
    """Simple vector figure for role-story slides."""
    add_circle(slide, x + 0.42 * scale, y, 0.58 * scale, fill=fill, line=accent)
    add_rect(
        slide,
        x + 0.22 * scale,
        y + 0.64 * scale,
        0.98 * scale,
        1.18 * scale,
        fill=fill,
        line=accent,
    )
    add_line(
        slide,
        x + 0.24 * scale,
        y + 0.94 * scale,
        x - 0.12 * scale,
        y + 1.42 * scale,
        accent,
        2.2,
    )
    add_line(
        slide,
        x + 1.18 * scale,
        y + 0.94 * scale,
        x + 1.56 * scale,
        y + 1.42 * scale,
        accent,
        2.2,
    )
    add_line(
        slide,
        x + 0.52 * scale,
        y + 1.8 * scale,
        x + 0.2 * scale,
        y + 2.38 * scale,
        accent,
        2.2,
    )
    add_line(
        slide,
        x + 0.9 * scale,
        y + 1.8 * scale,
        x + 1.22 * scale,
        y + 2.38 * scale,
        accent,
        2.2,
    )


def make_key_visual(path: Path) -> None:
    width, height = 1920, 1080
    base_color = "#F4F8FC" if LIGHT_THEME else "#06111E"
    grid_a = "#D9E8F1" if LIGHT_THEME else "#0B2C3D"
    grid_b = "#E4EEF5" if LIGHT_THEME else "#0A2232"
    node_fill = "#FFFFFF" if LIGHT_THEME else "#0D2234"
    hub_fill = "#E8F1F7" if LIGHT_THEME else "#0B2032"
    foreground = "#10243A" if LIGHT_THEME else "#F5F9FC"
    base = Image.new("RGB", (width, height), base_color)
    draw = ImageDraw.Draw(base)
    for x in range(-height, width, 92):
        draw.line((x, 0, x + height, height), fill=grid_a, width=1)
    for x in range(0, width, 120):
        draw.line((x, 0, x, height), fill=grid_b, width=1)
    for y in range(0, height, 90):
        draw.line((0, y, width, y), fill=grid_b, width=1)

    glow = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    router = (960, 520)
    nodes = [(420, 290), (1490, 285), (1510, 805)]
    accents = ["#37D7FF", "#35E6B1", "#2F73FF"]
    for (nx, ny), accent in zip(nodes, accents):
        for offset in (-18, 0, 18):
            gd.line(
                (router[0], router[1] + offset, nx, ny + offset // 2),
                fill=accent + "A0",
                width=4,
            )
        gd.ellipse((nx - 90, ny - 90, nx + 90, ny + 90), outline=accent, width=5)
        gd.rounded_rectangle(
            (nx - 54, ny - 58, nx + 54, ny + 58),
            radius=12,
            fill=node_fill,
            outline=accent,
            width=4,
        )
        for sy in range(ny - 38, ny + 40, 18):
            gd.line((nx - 32, sy, nx + 32, sy), fill=accent, width=2)
            gd.ellipse((nx + 20, sy - 3, nx + 27, sy + 4), fill=accent)

    gd.ellipse(
        (router[0] - 124, router[1] - 124, router[0] + 124, router[1] + 124),
        outline="#37D7FF",
        width=6,
    )
    gd.rounded_rectangle(
        (
            router[0] - 86,
            router[1] - 86,
            router[0] + 86,
            router[1] + 86,
        ),
        radius=24,
        fill=hub_fill,
        outline="#35E6B1",
        width=6,
    )
    gd.line(
        (router[0] - 48, router[1], router[0] + 48, router[1]),
        fill=foreground,
        width=8,
    )
    gd.polygon(
        [
            (router[0] + 48, router[1]),
            (router[0] + 22, router[1] - 18),
            (router[0] + 22, router[1] + 18),
        ],
        fill=foreground,
    )

    for _ in range(95):
        random.seed(_ * 731)
        x = random.randint(170, 1770)
        y = random.randint(130, 950)
        if math.dist((x, y), router) < 170:
            continue
        c = random.choice(["#37D7FF", "#35E6B1", "#2F73FF"])
        r = random.choice([2, 2, 3, 4])
        gd.ellipse((x - r, y - r, x + r, y + r), fill=c + "C0")

    blurred = glow.filter(ImageFilter.GaussianBlur(18))
    base = Image.alpha_composite(base.convert("RGBA"), blurred)
    base = Image.alpha_composite(base, glow)

    if not LIGHT_THEME:
        vignette = Image.new("L", (width, height), 0)
        vd = ImageDraw.Draw(vignette)
        for i in range(280):
            alpha = int((i / 280) ** 2 * 150)
            vd.rectangle((i, i, width - i, height - i), outline=alpha, width=3)
        black = Image.new("RGBA", (width, height), (0, 0, 0, 0))
        black.putalpha(vignette)
        base = Image.alpha_composite(base, black)
    base.convert("RGB").save(path, quality=95)


def make_tool_orbit_visual(path: Path) -> None:
    width, height = 1920, 1080
    base_color = "#F4F8FC" if LIGHT_THEME else "#06111E"
    grid_color = "#DDEAF2" if LIGHT_THEME else "#0B2638"
    orbit_fill = "#FFFFFF" if LIGHT_THEME else "#0A1D2ECC"
    chip_fill = "#E8F1F7" if LIGHT_THEME else "#102235"
    hub_fill = "#E8F1F7" if LIGHT_THEME else "#0B2032"
    foreground = "#10243A" if LIGHT_THEME else "#F5F9FC"
    base = Image.new("RGB", (width, height), base_color)
    grid = ImageDraw.Draw(base)
    for x in range(0, width, 96):
        grid.line((x, 0, x, height), fill=grid_color, width=1)
    for y in range(0, height, 96):
        grid.line((0, y, width, y), fill=grid_color, width=1)

    layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    draw = ImageDraw.Draw(layer)
    cx, cy = 960, 548
    labels = [
        ("OCR", "#37D7FF"),
        ("DOC", "#35E6B1"),
        ("BIOCLIP", "#2F73FF"),
        ("SAM", "#37D7FF"),
        ("YOLO", "#35E6B1"),
        ("TTS", "#FFB84A"),
        ("ASR", "#37D7FF"),
        ("LLM", "#2F73FF"),
        ("RAG", "#35E6B1"),
        ("ADDRESS", "#FFB84A"),
        ("WEB", "#37D7FF"),
        ("AUDIO", "#35E6B1"),
    ]
    title_font = ImageFont.truetype(
        "/home/john/.local/share/fonts/3wa-giday/Tektur-Regular.ttf", 44
    )
    label_font = ImageFont.truetype(
        "/home/john/.local/share/fonts/3wa-giday/Tektur-Regular.ttf", 25
    )
    hub_font = ImageFont.truetype(
        "/home/john/.local/share/fonts/3wa-giday/Tektur-Medium.ttf", 52
    )

    for i, (label, accent) in enumerate(labels):
        angle = -math.pi / 2 + i * (2 * math.pi / len(labels))
        radius_x = 690 if i % 2 == 0 else 615
        radius_y = 385 if i % 2 == 0 else 342
        nx = int(cx + math.cos(angle) * radius_x)
        ny = int(cy + math.sin(angle) * radius_y)
        draw.line((cx, cy, nx, ny), fill=accent + "78", width=4)
        draw.ellipse(
            (nx - 67, ny - 67, nx + 67, ny + 67),
            fill=orbit_fill,
            outline=accent,
            width=5,
        )
        draw.rounded_rectangle(
            (nx - 48, ny - 35, nx + 48, ny + 35),
            radius=12,
            fill=chip_fill,
            outline=accent,
            width=3,
        )
        bbox = draw.textbbox((0, 0), label, font=label_font)
        draw.text(
            (nx - (bbox[2] - bbox[0]) / 2, ny - (bbox[3] - bbox[1]) / 2 - 2),
            label,
            font=label_font,
            fill=foreground,
        )

    for radius, accent in [(190, "#37D7FF"), (150, "#35E6B1")]:
        draw.ellipse(
            (cx - radius, cy - radius, cx + radius, cy + radius),
            outline=accent,
            width=5,
        )
    draw.rounded_rectangle(
        (cx - 122, cy - 112, cx + 122, cy + 112),
        radius=30,
        fill=hub_fill,
        outline="#37D7FF",
        width=6,
    )
    hub = "AIHUB"
    bbox = draw.textbbox((0, 0), hub, font=hub_font)
    draw.text(
        (cx - (bbox[2] - bbox[0]) / 2, cy - 31),
        hub,
        font=hub_font,
        fill=foreground,
    )
    draw.text((86, 70), "CAPABILITY ORBIT", font=title_font, fill="#37D7FF")

    blurred = layer.filter(ImageFilter.GaussianBlur(18))
    base = Image.alpha_composite(base.convert("RGBA"), blurred)
    base = Image.alpha_composite(base, layer)
    base.convert("RGB").save(path, quality=95)


def build_deck() -> None:
    suffix = "-light" if LIGHT_THEME else ""
    key_visual = ASSETS / f"key-visual-sovereign-compute-mesh{suffix}.png"
    tool_visual = ASSETS / f"key-visual-aihub-tool-orbit{suffix}.png"
    make_key_visual(key_visual)
    make_tool_orbit_visual(tool_visual)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "GIS集團-AI服務平台｜GIDAY" + ("｜淡色版" if LIGHT_THEME else "")
    prs.core_properties.subject = "地端 GPU 服務治理、Cluster 與企業 AI 能力"
    prs.core_properties.author = "AI 算了吧..."
    prs.core_properties.keywords = "3waAIHub, GIS, GPU, Cluster, AI Governance"
    blank = prs.slide_layouts[6]

    # 01 Cover
    slide = prs.slides.add_slide(blank)
    add_image_crop(slide, key_visual, 0, 0, W, H)
    add_rect(slide, 0, 0, 7.9, H, fill=BG, radius=False, alpha=86)
    add_text(slide, "GIS GROUP × GIDAY 2026", 0.72, 0.62, 4.8, 0.3, 12, CYAN, True, DISPLAY)
    add_text(slide, "GIS集團－AI服務平台", 0.72, 1.22, 6.9, 0.7, 35, TEXT, True)
    add_runs(
        slide,
        [
            {"text": "把分散的 ", "color": MUTED, "size": 25},
            {"text": "GPU", "color": CYAN, "size": 25, "bold": True, "font": DISPLAY},
            {"text": "，變成可治理、可使用、可投資的企業能力", "color": TEXT, "size": 25, "bold": True},
        ],
        0.74,
        2.12,
        6.6,
        1.18,
    )
    add_rect(slide, 0.74, 3.58, 5.65, 0.03, fill=CYAN, radius=False)
    add_text(slide, "AI 算了吧...", 0.74, 4.02, 3.5, 0.38, 16, GREEN, True)
    add_text(slide, "08.07  ·  企業算力治理提案", 0.74, 4.48, 4.5, 0.38, 14, MUTED)
    add_text(
        slide,
        "買的是 GPU，累積的是企業能力。",
        0.74,
        6.44,
        5.6,
        0.42,
        18,
        TEXT,
        True,
    )
    add_text(slide, "01", 12.2, 7.02, 0.42, 0.24, 9, CYAN, True, DISPLAY, PP_ALIGN.RIGHT)

    # 02 I2C+S
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 2, "I²C+S 不只連結人，也應該連結算力", "WHY NOW")
    add_text(
        slide,
        "「一個人努力走不遠，\n但一群人努力，可以走得很遠。」",
        0.78,
        1.72,
        5.7,
        1.2,
        23,
        MUTED,
        True,
    )
    add_text(
        slide,
        "一台強的電腦也走不了很遠；\n但一群被治理、能協作的電腦，\n可以走得很遠。",
        0.78,
        3.03,
        5.9,
        1.9,
        27,
        TEXT,
        True,
    )
    add_rect(slide, 0.78, 5.4, 5.5, 0.04, fill=GREEN, radius=False)
    add_text(slide, "這就是 AIHub 的集團意義。", 0.78, 5.69, 5.2, 0.4, 18, GREEN, True)

    center_x, center_y = 9.7, 3.7
    add_circle(slide, center_x - 0.82, center_y - 0.82, 1.64, fill=BG2, line=CYAN)
    add_text(
        slide,
        "I²C+S",
        center_x - 0.76,
        center_y - 0.23,
        1.52,
        0.46,
        22,
        TEXT,
        True,
        DISPLAY,
        PP_ALIGN.CENTER,
        MSO_ANCHOR.MIDDLE,
    )
    concepts = [
        ("Integration", "主機＋技術＋服務", 7.0, 1.75, BLUE, "I"),
        ("Interoperability", "跨節點協同供應", 10.45, 1.75, CYAN, "I"),
        ("Communication", "狀態＋用量可見", 7.0, 4.95, AMBER, "C"),
        ("Sharing", "共享多餘算力", 10.45, 4.95, GREEN, "S"),
    ]
    for title, body, x, y, accent, mono in concepts:
        add_rect(slide, x, y, 2.2, 1.15, fill=PANEL, line=accent)
        add_text(slide, mono, x + 0.15, y + 0.13, 0.4, 0.35, 15, accent, True, DISPLAY)
        add_text(slide, title, x + 0.58, y + 0.14, 1.5, 0.28, 12, TEXT, True, DISPLAY)
        add_text(slide, body, x + 0.15, y + 0.59, 1.9, 0.28, 11, MUTED)
        add_line(slide, center_x, center_y, x + 1.1, y + 0.58, accent, 1.3)
    add_footer(slide, 2)

    # 03 Problem
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 3, "買到 AI 主機，不等於建立 AI 能力", "THE GAP")
    problems = [
        ("01", "專案各自採購", "規格、環境與責任人分散"),
        ("02", "服務各自安裝", "同一能力反覆建置"),
        ("03", "權限散落", "Token 與可用範圍難追查"),
        ("04", "GPU 同時閒置與壅塞", "需求看不到彼此"),
    ]
    xs = [0.74, 3.86, 6.98, 10.10]
    for index, (num, title, body) in enumerate(problems):
        x = xs[index]
        add_rect(slide, x, 1.74, 2.5, 2.42, fill=PANEL, line=RED if index == 3 else LINE)
        add_text(slide, num, x + 0.18, 1.94, 0.5, 0.3, 11, RED, True, DISPLAY)
        add_text(slide, title, x + 0.18, 2.47, 2.15, 0.58, 18, TEXT, True)
        add_text(slide, body, x + 0.18, 3.27, 2.0, 0.55, 12, MUTED)
        if index < 3:
            add_chevron(slide, x + 2.64, 2.62, 0.3, 0.52, RED)
    add_rect(slide, 0.74, 4.78, 11.86, 1.48, fill=RISK_PANEL, line=RED)
    stages = [
        ("硬體投資", TEXT),
        ("專案孤島", AMBER),
        ("維運風險", RED),
        ("ROI 看不見", RED),
    ]
    sx = 1.22
    for i, (label, color) in enumerate(stages):
        add_text(slide, label, sx, 5.21, 2.08, 0.42, 18, color, True, align=PP_ALIGN.CENTER)
        if i < len(stages) - 1:
            add_chevron(slide, sx + 2.18, 5.12, 0.38, 0.58, RED)
        sx += 2.9
    add_footer(slide, 3, "問題不是 GPU 不夠，而是缺少把 GPU 變成服務的治理層。")

    # 04 Platform
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 4, "AIHub 把主機投資，轉成可治理的服務供應鏈", "THE PRODUCT")
    add_image_crop(slide, ASSETS / "aihub-hubpack-market.png", 8.45, 1.55, 4.18, 4.9)
    add_rect(slide, 8.45, 1.55, 4.18, 4.9, fill=BG, radius=False, alpha=18)
    pipeline = [
        ("AI 主機", "GPU / CPU", BLUE),
        ("HubPack", "標準安裝", CYAN),
        ("Service Mode", "統一介面", TEAL),
        ("Account / Token", "授權範圍", AMBER),
        ("Router / Logs", "調度與追查", GREEN),
    ]
    y = 1.58
    for i, (title, body, accent) in enumerate(pipeline):
        add_rect(slide, 0.76, y, 6.85, 0.78, fill=PANEL, line=accent)
        add_text(slide, f"{i + 1:02d}", 0.95, y + 0.2, 0.5, 0.25, 11, accent, True, DISPLAY)
        add_text(slide, title, 1.55, y + 0.16, 2.15, 0.3, 16, TEXT, True, DISPLAY if i != 0 else CJK)
        add_text(slide, body, 4.15, y + 0.18, 2.8, 0.3, 14, MUTED)
        if i < len(pipeline) - 1:
            add_line(slide, 1.18, y + 0.78, 1.18, y + 0.96, accent, 1.5)
        y += 0.96
    add_rect(slide, 8.66, 5.42, 3.72, 1.02, fill=BG2, line=CYAN, alpha=94)
    add_text(slide, "18", 8.86, 5.59, 0.8, 0.5, 27, CYAN, True, DISPLAY)
    add_text(slide, "個 HubPack 目錄", 9.69, 5.66, 2.35, 0.32, 14, TEXT, True)
    add_text(slide, "安裝、健康檢查、API 文件與 Benchmark。", 8.86, 6.08, 3.15, 0.25, 9, MUTED)
    add_footer(slide, 4, "實機快照：3waAIHub HubPack 市集，2026-07-30。")

    # 05 Proof
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 5, "不是 PoC：節點已連線，服務已被呼叫", "PROOF")
    add_image_crop(slide, ASSETS / "aihub-dashboard-focusit.png", 5.65, 1.56, 6.95, 5.18)
    add_rect(slide, 5.65, 1.56, 6.95, 5.18, fill=BG, radius=False, alpha=8)
    add_metric(slide, 0.78, 1.68, 2.1, "3", "啟用節點", CYAN, "內網＋nature＋3wa.tw")
    add_metric(slide, 3.12, 1.68, 2.1, "14", "已安裝服務", GREEN, "本機服務清單")
    add_metric(slide, 0.78, 3.54, 2.1, "1,037", "API 呼叫", TEXT, "已寫入存取記錄")
    add_metric(slide, 3.12, 3.54, 2.1, "19", "已觀測 Mode", AMBER, "跨服務介面")
    add_rect(slide, 0.78, 5.46, 4.35, 0.98, fill=PANEL, line=CYAN)
    add_text(slide, "今天已能看到", 1.0, 5.68, 1.3, 0.26, 12, MUTED, True)
    add_text(slide, "節點健康、VRAM、服務與工作狀態", 2.32, 5.65, 2.55, 0.4, 14, TEXT, True)
    add_footer(slide, 5, "資料來源：統一入口 SQLite 與主控台；快照 2026-07-30。")

    # 06 Topology
    slide = prs.slides.add_slide(blank)
    add_image_crop(slide, ASSETS / "reference-compute-mesh.jpg", 0, 0, W, H)
    add_rect(slide, 0, 0, W, H, fill=BG, radius=False, alpha=80)
    add_header(slide, 6, "1 個統一入口，讓 3 台主機成為一個服務面", "DEPLOYMENT")
    router = add_rect(slide, 0.78, 2.73, 2.4, 1.55, fill=BG2, line=GREEN)
    add_text(slide, "UNIFIED ENTRY", 1.0, 2.95, 1.95, 0.25, 11, GREEN, True, DISPLAY, PP_ALIGN.CENTER)
    add_text(slide, "3wa.tw Router", 0.98, 3.34, 2.0, 0.36, 18, TEXT, True, align=PP_ALIGN.CENTER)
    add_text(slide, "單一 API 入口", 1.0, 3.78, 1.96, 0.25, 11, MUTED, align=PP_ALIGN.CENTER)
    nodes = [
        (4.35, 1.55, "3waAIHub Local", ["內網節點", "GPU 尚未回報", "2 modes"], AMBER, "A"),
        (7.73, 2.95, "Focusit AIHub Local", ["RTX 5090", "32,607 MB VRAM", "11 modes"], CYAN, "B"),
        (10.20, 4.55, "3wa.tw 服務節點", ["RTX 5060 Ti", "16,311 MB VRAM", "7 modes"], GREEN, "C"),
    ]
    for x, y, title, lines, accent, mono in nodes:
        add_node(slide, x, y, 2.58, 1.38, title, lines, accent, mono)
        add_line(slide, 3.18, 3.5, x, y + 0.68, accent, 2.2)
        add_chevron(slide, x - 0.20, y + 0.48, 0.25, 0.38, accent)
    add_text(slide, "身份驗證", 3.56, 4.48, 1.0, 0.3, 11, MUTED, True)
    add_text(slide, "Mode 路由", 5.1, 4.95, 1.0, 0.3, 11, MUTED, True)
    add_text(slide, "健康＋VRAM", 6.55, 5.42, 1.2, 0.3, 11, MUTED, True)
    add_text(slide, "用量紀錄", 8.3, 5.92, 1.0, 0.3, 11, MUTED, True)
    add_footer(slide, 6, "拓樸為統一入口視角；VRAM 與 Mode 取自 2026-07-30 節點快照。")

    # 07 Governance
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 7, "真正的治理，發生在每一次 API 呼叫", "GOVERNANCE")
    steps = [
        ("WHO", "帳號", "誰在使用", BLUE),
        ("KEY", "Token", "能否進入", CYAN),
        ("WHAT", "Mode 權限", "能用什麼", TEAL),
        ("WHERE", "節點路由", "派去哪裡", AMBER),
        ("PROOF", "Logs / Usage", "結果與成本", GREEN),
    ]
    x = 0.76
    for i, (tag, title, body, accent) in enumerate(steps):
        add_rect(slide, x, 1.73, 2.08, 1.46, fill=PANEL, line=accent)
        add_text(slide, tag, x + 0.18, 1.92, 0.78, 0.23, 9, accent, True, DISPLAY)
        add_text(slide, title, x + 0.18, 2.25, 1.74, 0.32, 16, TEXT, True)
        add_text(slide, body, x + 0.18, 2.71, 1.62, 0.25, 11, MUTED)
        if i < len(steps) - 1:
            add_chevron(slide, x + 2.16, 2.16, 0.28, 0.48, accent)
        x += 2.5
    add_image_crop(slide, ASSETS / "aihub-cluster-management.png", 0.78, 3.6, 6.45, 2.82)
    add_rect(slide, 7.56, 3.6, 5.04, 2.82, fill=PANEL, line=LINE)
    qs = [
        ("誰？", "member_id / token_id"),
        ("能用什麼？", "service permission / mode"),
        ("派去哪裡？", "station / health / VRAM"),
        ("結果如何？", "status / elapsed / bytes"),
    ]
    y = 3.88
    for q, a in qs:
        add_text(slide, q, 7.86, y, 1.2, 0.3, 13, CYAN, True)
        add_text(slide, a, 9.22, y, 2.95, 0.3, 12, MUTED, font=DISPLAY)
        y += 0.56
    add_footer(slide, 7, "實機截圖中的 Cluster Token 已遮蔽。")

    # 08 VRAM
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 8, "有 GPU，不等於有可用 VRAM", "CAPACITY")
    add_text(slide, "目前節點快照", 0.78, 1.55, 3.1, 0.35, 15, TEXT, True)
    hosts = [
        ("Focusit / RTX 5090", 32.607, 32.2, RED, "98.8% 已用"),
        ("3wa.tw / RTX 5060 Ti", 16.311, 5.825, AMBER, "35.7% 已用"),
        ("內網節點", 0, 0, MUTED, "GPU 尚未回報"),
    ]
    y = 2.02
    max_gb = 36.0
    for name, total, used, accent, note in hosts:
        add_text(slide, name, 0.78, y, 2.25, 0.3, 12, TEXT, True)
        add_rect(slide, 3.04, y + 0.03, 3.35, 0.25, fill=BAR_BG, radius=False)
        if total:
            add_rect(
                slide,
                3.04,
                y + 0.03,
                3.35 * used / max_gb,
                0.25,
                fill=accent,
                radius=False,
            )
            add_line(
                slide,
                3.04 + 3.35 * total / max_gb,
                y - 0.01,
                3.04 + 3.35 * total / max_gb,
                y + 0.33,
                CYAN,
                1,
            )
        add_text(slide, note, 6.55, y - 0.02, 1.35, 0.3, 11, accent, True, align=PP_ALIGN.RIGHT)
        y += 0.66
    add_text(slide, "HubPack 宣告最低 VRAM", 0.78, 4.28, 3.4, 0.35, 15, TEXT, True)
    packs = [
        ("GPU Audio Cleanup", 8.192, BLUE),
        ("VoxCPM2 TTS", 9.6, CYAN),
        ("TranslateGemma 12B", 10.0, TEAL),
        ("SAM3", 16.0, AMBER),
        ("Gemma 4 12B", 16.0, RED),
    ]
    x0 = 0.78
    y0 = 4.86
    bar_w = 7.1
    for i, (name, value, accent) in enumerate(packs):
        yy = y0 + i * 0.36
        add_text(slide, name, x0, yy - 0.03, 2.05, 0.24, 10, MUTED)
        add_rect(slide, x0 + 2.05, yy, bar_w * value / 32, 0.18, fill=accent, radius=False)
        add_text(slide, f"{value:g} GB", x0 + 2.14 + bar_w * value / 32, yy - 0.06, 0.68, 0.25, 10, accent, True, font=DISPLAY)
    add_rect(slide, 8.72, 1.66, 3.9, 4.9, fill=PANEL, line=LINE)
    add_text(slide, "這張快照說了什麼？", 9.02, 1.99, 3.1, 0.42, 20, TEXT, True)
    add_text(slide, "32 GB", 9.02, 2.75, 1.3, 0.46, 27, RED, True, DISPLAY)
    add_text(slide, "節點近滿載", 10.25, 2.82, 1.6, 0.3, 14, TEXT, True)
    add_text(slide, "16 GB", 9.02, 3.55, 1.3, 0.46, 27, AMBER, True, DISPLAY)
    add_text(slide, "節點仍有空間", 10.25, 3.62, 1.7, 0.3, 14, TEXT, True)
    add_rect(slide, 9.02, 4.36, 3.05, 0.03, fill=CYAN, radius=False)
    add_text(
        slide,
        "缺的不是「更多卡」這個單一答案，\n而是先看見容量、再做路由與採購。",
        9.02,
        4.72,
        3.1,
        1.15,
        18,
        TEXT,
        True,
    )
    add_footer(
        slide,
        8,
        "HubPack 最低 VRAM ≠ 實測峰值；實際用量受模型、量化、批次、上下文與併發影響。",
    )

    # 09 Balance
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 9, "不是每台都買最大，是讓每張卡都被看見", "BALANCE")
    add_text(slide, "今天的快照", 0.78, 1.56, 4.3, 0.35, 15, RED, True)
    add_text(slide, "治理後的目標", 7.0, 1.56, 4.3, 0.35, 15, GREEN, True)
    names = [("內網節點", 0, MUTED), ("RTX 5090", 99, RED), ("RTX 5060 Ti", 36, AMBER)]
    y = 2.14
    for name, value, accent in names:
        add_text(slide, name, 0.78, y, 1.55, 0.26, 12, TEXT, True)
        add_rect(slide, 2.35, y + 0.02, 3.0, 0.28, fill=BAR_BG, radius=False)
        if value:
            add_rect(slide, 2.35, y + 0.02, 3.0 * value / 100, 0.28, fill=accent, radius=False)
        add_text(slide, f"{value if value else 'N/A'}{'%' if value else ''}", 5.52, y - 0.02, 0.55, 0.3, 11, accent, True, font=DISPLAY, align=PP_ALIGN.RIGHT)
        y += 0.66
    add_text(
        slide,
        "同時存在：\n一台近滿載、一台有餘裕、一台沒有資料。",
        0.78,
        4.35,
        5.2,
        1.0,
        19,
        TEXT,
        True,
    )
    add_chevron(slide, 6.17, 3.03, 0.55, 0.9, CYAN)
    add_rect(slide, 7.0, 2.02, 5.56, 3.4, fill=PANEL, line=GREEN)
    target = [
        ("1", "依 Mode 找得到服務", CYAN),
        ("2", "依健康與 VRAM 選節點", TEAL),
        ("3", "排隊、超時、失敗都有紀錄", AMBER),
        ("4", "讓採購決策有實際使用證據", GREEN),
    ]
    yy = 2.34
    for num, text, accent in target:
        add_circle(slide, 7.3, yy, 0.46, fill=BG2, line=accent)
        add_text(slide, num, 7.3, yy, 0.46, 0.46, 12, accent, True, DISPLAY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(slide, text, 7.95, yy + 0.04, 4.1, 0.34, 16, TEXT, True)
        yy += 0.72
    outcomes = [("降低等待", CYAN), ("提高共用", GREEN), ("延後重複採購", AMBER)]
    xx = 7.05
    for label, accent in outcomes:
        add_rect(slide, xx, 5.76, 1.68, 0.57, fill=BG2, line=accent)
        add_text(slide, label, xx, 5.89, 1.68, 0.26, 12, accent, True, align=PP_ALIGN.CENTER)
        xx += 1.86
    add_footer(slide, 9, "負載平衡是治理目標；目前快照僅代表當下容量狀態。")

    # 10 Audience journey
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 10, "不同角色，從同一平台拿到不同答案", "AUDIENCE")
    roles = [
        ("工程師", "一個 endpoint\n一個 Token\n可複製範例", CYAN, "ENG"),
        ("SA／規劃師", "節點、Pack、VRAM\n可重用架構", BLUE, "SA"),
        ("PM", "呼叫、成功率、容量\n可說明範圍與成本", AMBER, "PM"),
        ("主管／外賓", "資產利用、服務能力\n可看見投資成果", GREEN, "ROI"),
    ]
    x = 0.72
    for role, body, accent, mono in roles:
        add_rect(slide, x, 1.62, 2.88, 2.18, fill=PANEL, line=accent)
        add_text(slide, mono, x + 0.2, 1.84, 0.74, 0.27, 10, accent, True, DISPLAY)
        add_text(slide, role, x + 0.2, 2.24, 2.36, 0.38, 18, TEXT, True)
        add_text(slide, body, x + 0.2, 2.78, 2.35, 0.8, 13, MUTED)
        x += 3.13
    add_text(slide, "現場 Demo 路線", 0.75, 4.35, 2.3, 0.35, 17, TEXT, True)
    demos = [
        ("DEMO 1", "Token → Edge TTS", "證明工程師能立即介接", CYAN),
        ("DEMO 2", "文件／地址／影像", "證明能力能跨專案重用", TEAL),
        ("DEMO 3", "Agent → AIHub Tools", "證明未來能由意圖調度", GREEN),
    ]
    x = 0.75
    for tag, title, body, accent in demos:
        add_rect(slide, x, 4.88, 3.82, 1.38, fill=BG2, line=accent)
        add_text(slide, tag, x + 0.2, 5.08, 0.9, 0.25, 9, accent, True, DISPLAY)
        add_text(slide, title, x + 1.15, 5.03, 2.35, 0.35, 15, TEXT, True, DISPLAY if "Token" in title or "Agent" in title else CJK)
        add_text(slide, body, x + 0.2, 5.59, 3.25, 0.3, 12, MUTED)
        x += 4.05
    add_footer(slide, 10, "每個 Demo 都回到同一條治理鏈：身份、權限、路由、紀錄。")

    # 11 Demo API
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 11, "Demo 1｜拿到 Token，工程師就能呼叫服務", "LIVE API")
    add_image_crop(slide, ASSETS / "aihub-api-playground-edge-tts.png", 0.72, 1.54, 7.55, 4.98)
    add_rect(slide, 8.58, 1.54, 4.02, 4.98, fill=PANEL, line=LINE)
    demo_steps = [
        ("01", "選擇 mode=edge_tts", "同一入口也能切換 OCR、地址、YOLO、chat"),
        ("02", "調整參數", "voice / rate / volume / pitch / subtitles"),
        ("03", "複製程式碼", "curl / PHP / JavaScript 直接介接"),
    ]
    y = 1.86
    for num, title, body in demo_steps:
        add_text(slide, num, 8.88, y, 0.55, 0.28, 10, CYAN, True, DISPLAY)
        add_text(slide, title, 9.55, y - 0.02, 2.6, 0.34, 16, TEXT, True)
        add_text(slide, body, 9.55, y + 0.38, 2.53, 0.52, 11, MUTED)
        y += 1.15
    add_rect(slide, 8.88, 5.28, 3.4, 0.78, fill=BG2, line=GREEN)
    add_text(slide, "同一呼叫模式，跨越不同服務與節點。", 9.08, 5.48, 3.0, 0.38, 13, GREEN, True)
    add_footer(slide, 11, "畫面中的 Bearer Token 為示意字串，不包含真實金鑰。")

    # 12 Demo workflows
    slide = prs.slides.add_slide(blank)
    add_image_crop(slide, ASSETS / "reference-document-workflow.jpg", 0, 0, W, H)
    add_rect(slide, 0, 0, W, H, fill=BG, radius=False, alpha=76)
    add_header(slide, 12, "Demo 2｜專案功能，變成可重用的服務鏈", "REUSE")
    flows = [
        ("A", "技術文件", "PDF → DocParser / OCR → Markdown / 圖片 → RAG", CYAN),
        ("B", "台灣地址", "混亂門牌 → Wash / Geocode → GIS 專案資料", GREEN),
        ("C", "影像任務", "Image → 去背 / YOLO / BioCLIP / SAM3 → 空間流程", AMBER),
    ]
    y = 1.67
    for mono, title, body, accent in flows:
        add_rect(slide, 0.76, y, 6.55, 1.26, fill=BG2, line=accent, alpha=94)
        add_circle(slide, 1.0, y + 0.26, 0.64, fill=BG, line=accent)
        add_text(slide, mono, 1.0, y + 0.26, 0.64, 0.64, 16, accent, True, DISPLAY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(slide, title, 1.86, y + 0.22, 1.2, 0.35, 17, TEXT, True)
        add_text(slide, body, 3.15, y + 0.22, 3.74, 0.68, 13, MUTED)
        y += 1.44
    add_rect(slide, 7.65, 1.68, 4.95, 4.15, fill=BG2, line=LINE, alpha=92)
    add_text(slide, "重用的不是一個模型，\n而是一條可治理的能力鏈。", 8.0, 2.12, 4.2, 1.05, 24, TEXT, True)
    add_rect(slide, 8.0, 3.58, 3.85, 0.04, fill=CYAN, radius=False)
    add_text(slide, "一次標準化", 8.0, 4.02, 1.75, 0.35, 16, CYAN, True)
    add_text(slide, "多專案呼叫", 10.02, 4.02, 1.75, 0.35, 16, GREEN, True)
    add_text(slide, "同一套 Token、路由、工作狀態與 artifacts。", 8.0, 4.72, 3.95, 0.7, 15, MUTED)
    add_footer(slide, 12, "文件工作流圖為概念視覺；Demo 以 3waAIHub 實際服務為準。")

    # 13 Usage & chargeback
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 13, "有 Token 與用量紀錄，才有共享與計價的底座", "ECONOMICS")
    add_image_crop(slide, ASSETS / "aihub-cluster-usage-masked.png", 0.72, 1.55, 7.55, 4.92)
    add_rect(slide, 8.55, 1.55, 4.04, 3.56, fill=PANEL, line=LINE)
    dims = [
        ("WHO", "帳號／Token", CYAN),
        ("WHAT", "Mode／服務", TEAL),
        ("WHERE", "子節點", BLUE),
        ("RESULT", "成功／失敗／容量", AMBER),
    ]
    y = 1.89
    for tag, label, accent in dims:
        add_text(slide, tag, 8.86, y, 0.85, 0.23, 9, accent, True, DISPLAY)
        add_text(slide, label, 9.92, y - 0.03, 2.25, 0.3, 15, TEXT, True)
        y += 0.68
    outcomes = [
        ("配額", "Quota"),
        ("預算", "Budget"),
        ("內部分攤", "Chargeback"),
        ("對外服務", "Service"),
    ]
    x = 8.55
    for i, (zh, en) in enumerate(outcomes):
        xx = x + (i % 2) * 2.05
        yy = 5.42 + (i // 2) * 0.67
        add_rect(slide, xx, yy, 1.86, 0.52, fill=BG2, line=GREEN if i >= 2 else CYAN)
        add_text(slide, zh, xx + 0.15, yy + 0.11, 0.8, 0.24, 12, TEXT, True)
        add_text(slide, en, xx + 0.9, yy + 0.12, 0.75, 0.22, 8, MUTED, True, DISPLAY, PP_ALIGN.RIGHT)
    add_footer(slide, 13, "實機用量畫面已遮蔽會員與 Token；計價流程為治理機制示意。")

    # 14 Agent + roadmap
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 14, "Agent 會把 AIHub 變成可行動的企業能力層", "NEXT")
    columns = [
        (0.72, 2.02, 2.0, "使用者意圖", ["Hermes", "OpenAI Agents", "其他 Agent"], BLUE),
        (3.02, 2.02, 2.0, "治理與核准", ["身分／Token", "Allowed tools", "敏感動作核准"], AMBER),
        (5.32, 2.02, 2.0, "AIHub Cluster", ["統一 API", "Mode 路由", "節點健康／VRAM"], CYAN),
        (7.62, 2.02, 2.0, "內部工具", ["Doc／OCR", "Address／Vision", "TTS／LLM"], TEAL),
        (9.92, 2.02, 2.68, "Trace & Usage", ["工具呼叫", "結果／錯誤", "用量／成本"], GREEN),
    ]
    for i, (x, y, w, title, lines, accent) in enumerate(columns):
        add_rect(slide, x, y, w, 1.86, fill=PANEL, line=accent)
        add_text(slide, f"{i + 1:02d}", x + 0.16, y + 0.18, 0.38, 0.22, 9, accent, True, DISPLAY)
        add_text(slide, title, x + 0.16, y + 0.55, w - 0.32, 0.34, 16, TEXT, True)
        add_text(slide, "\n".join(lines), x + 0.16, y + 1.02, w - 0.32, 0.65, 11, MUTED)
        if i < len(columns) - 1:
            add_chevron(slide, x + w + 0.06, y + 0.68, 0.22, 0.45, accent)
    add_text(slide, "90 天落地路線", 0.72, 4.35, 2.0, 0.35, 16, TEXT, True)
    roadmap = [
        ("0–30", "盤點與治理", "節點／Token／健康／用量", BLUE),
        ("31–60", "標準化與 Demo", "HubPack／範例／權限", CYAN),
        ("61–90", "Agent 與計價試點", "Tools／核准／Chargeback", GREEN),
    ]
    x = 0.72
    for phase, title, body, accent in roadmap:
        add_rect(slide, x, 4.87, 3.86, 1.12, fill=BG2, line=accent)
        add_text(slide, phase, x + 0.18, 5.06, 0.8, 0.27, 11, accent, True, DISPLAY)
        add_text(slide, title, x + 1.08, 5.03, 2.45, 0.3, 15, TEXT, True)
        add_text(slide, body, x + 1.08, 5.43, 2.45, 0.26, 10, MUTED)
        x += 4.02
    add_text(
        slide,
        "OpenAI 官方文件：Agents SDK 支援 tools、handoffs、guardrails、tracing；Responses API 可連接 remote MCP，敏感工具可要求核准。",
        0.72,
        6.34,
        11.5,
        0.46,
        9,
        MUTED,
    )
    add_footer(slide, 14, "來源：developers.openai.com/api/docs/guides/agents、tools-connectors-mcp；存取 2026-07-30。")

    # 15 Chapter transition
    slide = prs.slides.add_slide(blank)
    add_image_crop(slide, tool_visual, 0, 0, W, H)
    add_rect(slide, 0, 0, W, H, fill=BG, radius=False, alpha=45)
    add_rect(slide, 0, 0, 7.25, H, fill=BG, radius=False, alpha=84)
    add_text(slide, "15  INSIDE THE PRODUCT", 0.72, 0.58, 3.6, 0.28, 11, CYAN, True, DISPLAY)
    add_text(slide, "安裝 AIHub 之後，\n一台主機開始長出手腳", 0.72, 1.38, 6.25, 1.65, 34, TEXT, True)
    add_text(
        slide,
        "Market 提供能力，HubPack 標準化安裝，\nAPI 讓工程師與 Agent 快速組合。",
        0.74,
        3.45,
        5.9,
        0.92,
        19,
        MUTED,
        True,
    )
    add_text(slide, "接下來，拆開系統看。", 0.74, 5.63, 4.6, 0.42, 20, GREEN, True)
    add_text(slide, "15", 12.2, 7.02, 0.42, 0.24, 9, CYAN, True, DISPLAY, PP_ALIGN.RIGHT)

    # 16 Capability orbit
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 16, "AIHub 讓硬體一次獲得多種可組合能力", "CAPABILITY")
    add_image_crop(slide, tool_visual, 0.72, 1.52, 8.0, 4.94)
    add_rect(slide, 8.98, 1.52, 3.63, 1.27, fill=PANEL, line=CYAN)
    add_text(slide, "看", 9.22, 1.79, 0.55, 0.42, 22, CYAN, True)
    add_text(slide, "OCR／YOLO／SAM／BiRefNet", 9.92, 1.83, 2.42, 0.32, 12, TEXT, True)
    add_rect(slide, 8.98, 2.98, 3.63, 1.27, fill=PANEL, line=TEAL)
    add_text(slide, "聽說", 9.22, 3.25, 0.78, 0.42, 22, TEAL, True)
    add_text(slide, "ASR／TTS／Audio Cleanup", 10.1, 3.29, 2.2, 0.32, 12, TEXT, True)
    add_rect(slide, 8.98, 4.44, 3.63, 1.27, fill=PANEL, line=AMBER)
    add_text(slide, "思考", 9.22, 4.71, 0.78, 0.42, 22, AMBER, True)
    add_text(slide, "LLM／RAG／翻譯／地址清洗", 10.1, 4.75, 2.2, 0.32, 12, TEXT, True)
    add_text(slide, "不是裝一個模型，是取得一組可治理的數位手腳。", 8.98, 6.08, 3.62, 0.46, 14, GREEN, True)
    add_footer(slide, 16, "依目前 18 個 HubPack 能力分類整理；單一節點可按硬體條件選擇安裝。")

    # 17 Market
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 17, "Market 是企業 AI 能力目錄，不只是模型倉庫", "MARKET")
    add_image_crop(slide, ASSETS / "aihub-hubpack-market.png", 0.72, 1.53, 6.55, 4.95)
    market_groups = [
        ("SEE", "看懂影像", "SAM3 · YOLO · OCR · BiRefNet", CYAN),
        ("READ", "讀懂文件", "Docparser · PPStructure · RAG", TEAL),
        ("TALK", "聽與說", "Whisper ASR · Edge TTS · VoxCPM", AMBER),
        ("ACT", "完成任務", "地址清洗 · 網頁截圖 · 音訊清理", GREEN),
    ]
    y = 1.55
    for tag, title, examples, accent in market_groups:
        add_rect(slide, 7.58, y, 5.03, 1.05, fill=PANEL, line=accent)
        add_text(slide, tag, 7.82, y + 0.2, 0.68, 0.26, 10, accent, True, DISPLAY)
        add_text(slide, title, 8.64, y + 0.17, 1.2, 0.32, 16, TEXT, True)
        add_text(slide, examples, 8.64, y + 0.57, 3.62, 0.25, 10, MUTED)
        y += 1.23
    add_text(slide, "18", 7.65, 6.24, 1.0, 0.42, 24, CYAN, True, DISPLAY)
    add_text(slide, "個現有 Pack，可依節點條件逐步擴充", 8.58, 6.3, 3.72, 0.28, 13, TEXT, True)
    add_footer(slide, 17, "Market 實機畫面與本機 packs/ 目錄快照：2026-07-30。")

    # 18 Install flow
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 18, "從按下安裝，到可被 API 呼叫，是一條標準流程", "INSTALL")
    steps = [
        ("01", "選 Pack", "用途／硬體需求", BLUE),
        ("02", "Preflight", "Docker／VRAM／Port", CYAN),
        ("03", "安裝", "環境與相依套件", TEAL),
        ("04", "設定", "URL／模型／參數", AMBER),
        ("05", "啟動", "Service runtime", BLUE),
        ("06", "健康檢查", "Ready 才啟用", GREEN),
        ("07", "API 測試", "範例與回應", CYAN),
        ("08", "Benchmark", "容量與品質", GREEN),
    ]
    positions = [
        (0.74, 1.7), (3.76, 1.7), (6.78, 1.7), (9.8, 1.7),
        (9.8, 4.4), (6.78, 4.4), (3.76, 4.4), (0.74, 4.4),
    ]
    for i, ((num, title, body, accent), (x, y)) in enumerate(zip(steps, positions)):
        add_rect(slide, x, y, 2.52, 1.42, fill=PANEL, line=accent)
        add_text(slide, num, x + 0.18, y + 0.2, 0.44, 0.26, 10, accent, True, DISPLAY)
        add_text(slide, title, x + 0.18, y + 0.57, 2.0, 0.34, 17, TEXT, True)
        add_text(slide, body, x + 0.18, y + 1.03, 2.0, 0.22, 10, MUTED)
        if i < 3:
            add_chevron(slide, x + 2.65, y + 0.46, 0.27, 0.48, accent)
        if 4 <= i < 7:
            add_chevron(slide, x - 0.38, y + 0.46, 0.27, 0.48, accent)
    add_line(slide, 11.06, 3.15, 11.06, 4.35, AMBER, 2)
    add_text(slide, "失敗時停在該步驟，留下可追查紀錄", 0.76, 6.33, 5.1, 0.38, 15, AMBER, True)
    add_text(slide, "成功後才成為可授權、可測試的服務", 7.43, 6.33, 5.1, 0.38, 15, GREEN, True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 18, "流程依 AIHub 現有安裝、設定、啟動、健康檢查與 Benchmark 操作整理。")

    # 19 Pack anatomy
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 19, "一個 HubPack，把重複研發工作包成共同規格", "PACK")
    add_circle(slide, 5.52, 2.45, 2.25, fill=BG2, line=CYAN)
    add_text(slide, "HubPack", 5.52, 3.05, 2.25, 0.42, 24, TEXT, True, DISPLAY, PP_ALIGN.CENTER)
    add_text(slide, "可安裝能力單元", 5.52, 3.56, 2.25, 0.3, 12, MUTED, True, align=PP_ALIGN.CENTER)
    pack_parts = [
        ("用途與版本", "pack.json", 0.78, 1.62, BLUE),
        ("硬體門檻", "GPU／VRAM／Disk", 0.78, 4.57, AMBER),
        ("執行環境", "Docker／CLI／Service", 4.23, 5.15, TEAL),
        ("健康契約", "ready／timeout／error", 8.35, 5.15, GREEN),
        ("API 契約", "mode／input／output", 10.0, 1.62, CYAN),
        ("維運工具", "logs／rebuild／benchmark", 8.75, 3.36, BLUE),
    ]
    for title, body, x, y, accent in pack_parts:
        add_rect(slide, x, y, 2.55, 1.05, fill=PANEL, line=accent)
        add_text(slide, title, x + 0.17, y + 0.17, 2.15, 0.28, 15, TEXT, True)
        add_text(slide, body, x + 0.17, y + 0.6, 2.15, 0.24, 10, accent, True, DISPLAY)
        add_line(slide, 6.65, 3.56, x + 1.28, y + 0.53, accent, 1.3)
    add_footer(slide, 19, "標準化的價值：一次解決安裝、介面、檢查與維運，不讓每個專案重做。")

    # 20 Service lifecycle
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 20, "工具上線之後，仍要被持續治理", "OPERATIONS")
    lifecycle = [
        ("INSTALL", "安裝", "版本與相依", BLUE),
        ("RUN", "執行", "服務狀態", CYAN),
        ("OBSERVE", "觀測", "Logs／Usage", TEAL),
        ("UPDATE", "更新", "Pack／模型", AMBER),
        ("REBUILD", "重建", "環境復原", BLUE),
        ("RETIRE", "退場", "停止與封存", RED),
    ]
    x = 0.76
    for i, (tag, title, body, accent) in enumerate(lifecycle):
        add_circle(slide, x, 2.08, 1.37, fill=BG2, line=accent)
        add_text(slide, f"{i + 1:02d}", x, 2.38, 1.37, 0.28, 11, accent, True, DISPLAY, PP_ALIGN.CENTER)
        add_text(slide, title, x - 0.12, 3.67, 1.62, 0.35, 17, TEXT, True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x - 0.22, 4.14, 1.82, 0.32, 11, MUTED, align=PP_ALIGN.CENTER)
        add_text(slide, tag, x - 0.12, 4.65, 1.62, 0.25, 9, accent, True, DISPLAY, PP_ALIGN.CENTER)
        if i < len(lifecycle) - 1:
            add_chevron(slide, x + 1.52, 2.49, 0.27, 0.48, accent)
        x += 2.02
    add_rect(slide, 1.35, 5.55, 10.65, 0.72, fill=PANEL, line=GREEN)
    add_text(slide, "安裝不是終點；可觀測、可更新、可退場，才是企業平台。", 1.35, 5.74, 10.65, 0.32, 18, GREEN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 20)

    # 21 Architecture
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 21, "這不是一個工具頁，而是一張企業 AI 服務網路", "ARCHITECTURE")
    add_text(slide, "需求端", 0.78, 1.54, 1.2, 0.3, 11, CYAN, True, DISPLAY)
    actors = [("工程師", "DEV"), ("PM／SA", "PLAN"), ("系統", "APP"), ("Agent", "AGENT")]
    y = 1.95
    for title, mono in actors:
        add_rect(slide, 0.76, y, 1.65, 0.72, fill=PANEL, line=LINE)
        add_text(slide, mono, 0.92, y + 0.15, 0.52, 0.22, 8, CYAN, True, DISPLAY)
        add_text(slide, title, 1.38, y + 0.18, 0.88, 0.25, 12, TEXT, True)
        add_line(slide, 2.41, y + 0.36, 3.02, 3.47, CYAN, 1.2)
        y += 0.91

    add_rect(slide, 3.02, 1.72, 2.16, 4.7, fill=BG2, line=BLUE)
    add_text(slide, "統一入口", 3.28, 1.98, 1.66, 0.36, 18, TEXT, True, align=PP_ALIGN.CENTER)
    gateway = [
        ("AUTH", "帳號／Token"),
        ("POLICY", "權限／Quota"),
        ("ROUTER", "Mode／節點"),
        ("APPROVAL", "敏感動作"),
        ("TRACE", "請求／結果"),
    ]
    gy = 2.59
    for tag, label in gateway:
        add_rect(slide, 3.31, gy, 1.58, 0.55, fill=PANEL, line=CYAN)
        add_text(slide, tag, 3.47, gy + 0.11, 0.65, 0.2, 8, CYAN, True, DISPLAY)
        add_text(slide, label, 4.06, gy + 0.11, 0.68, 0.22, 9, TEXT, True, align=PP_ALIGN.RIGHT)
        gy += 0.67

    add_text(slide, "算力供應端", 5.79, 1.54, 1.45, 0.3, 11, TEAL, True, DISPLAY)
    stations = [
        ("Station A", "3wa 主端點／GPU", 5.78, 1.95, CYAN),
        ("Station B", "Focusit／GPU", 5.78, 3.37, TEAL),
        ("Station C", "內網節點／GPU", 5.78, 4.79, BLUE),
    ]
    for title, body, x, y, accent in stations:
        add_rect(slide, x, y, 2.15, 1.0, fill=PANEL, line=accent)
        add_text(slide, title, x + 0.18, y + 0.17, 1.5, 0.26, 14, TEXT, True, font=DISPLAY)
        add_text(slide, body, x + 0.18, y + 0.59, 1.62, 0.23, 10, MUTED)
        add_line(slide, 5.18, 3.48, x, y + 0.5, accent, 1.4)

    add_text(slide, "服務層", 8.42, 1.54, 1.0, 0.3, 11, AMBER, True, DISPLAY)
    services = [("VISION", "OCR／SAM／YOLO"), ("LANGUAGE", "LLM／RAG／翻譯"), ("AUDIO", "ASR／TTS"), ("GIS", "地址／文件／網頁")]
    sy = 1.95
    for tag, label in services:
        add_rect(slide, 8.4, sy, 1.78, 0.82, fill=PANEL, line=AMBER)
        add_text(slide, tag, 8.56, sy + 0.13, 1.0, 0.2, 8, AMBER, True, DISPLAY)
        add_text(slide, label, 8.56, sy + 0.43, 1.38, 0.22, 10, TEXT, True)
        add_line(slide, 7.93, 3.45, 8.4, sy + 0.41, AMBER, 1.2)
        sy += 1.04

    add_text(slide, "治理資料", 10.72, 1.54, 1.2, 0.3, 11, GREEN, True, DISPLAY)
    data_nodes = [
        ("MARKET", "Pack／版本"),
        ("REGISTRY", "服務／端點"),
        ("USAGE", "Token／用量"),
        ("LOGS", "API／任務／系統"),
        ("METRICS", "健康／容量"),
    ]
    dy = 1.95
    for tag, label in data_nodes:
        add_rect(slide, 10.69, dy, 1.92, 0.7, fill=PANEL, line=GREEN)
        add_text(slide, tag, 10.85, dy + 0.13, 0.76, 0.2, 8, GREEN, True, DISPLAY)
        add_text(slide, label, 11.55, dy + 0.13, 0.88, 0.22, 9, TEXT, True, align=PP_ALIGN.RIGHT)
        add_line(slide, 10.18, 3.47, 10.69, dy + 0.35, GREEN, 1.1)
        dy += 0.85
    add_footer(slide, 21, "概念架構依 AIHub 現有後台、Cluster、Pack、Token、Router、Logs 與 Metrics 功能整理。")

    # 22 Request path
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 22, "一次 API 呼叫，背後經過七個治理檢查點", "REQUEST PATH")
    request_steps = [
        ("01", "帶 Token", "WHO", CYAN),
        ("02", "驗權限", "CAN", BLUE),
        ("03", "選 Mode", "WHAT", TEAL),
        ("04", "挑節點", "WHERE", AMBER),
        ("05", "跑服務", "DO", CYAN),
        ("06", "留紀錄", "TRACE", GREEN),
        ("07", "回結果", "RESULT", GREEN),
    ]
    x = 0.63
    for i, (num, title, tag, accent) in enumerate(request_steps):
        add_circle(slide, x, 2.2, 1.28, fill=BG2, line=accent)
        add_text(slide, num, x, 2.5, 1.28, 0.25, 10, accent, True, DISPLAY, PP_ALIGN.CENTER)
        add_text(slide, title, x - 0.18, 3.76, 1.64, 0.35, 17, TEXT, True, align=PP_ALIGN.CENTER)
        add_text(slide, tag, x - 0.1, 4.25, 1.48, 0.24, 9, accent, True, DISPLAY, PP_ALIGN.CENTER)
        if i < len(request_steps) - 1:
            add_chevron(slide, x + 1.39, 2.59, 0.25, 0.46, accent)
        x += 1.79
    add_rect(slide, 1.22, 5.25, 10.88, 0.86, fill=PANEL, line=GREEN)
    add_text(slide, "同一支 API，能被不同人、不同 Agent 使用；治理與稽核仍留在同一入口。", 1.22, 5.49, 10.88, 0.35, 17, GREEN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 22)

    # 23 Governance navigation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 23, "後台不是一頁，而是一個 AI 治理控制面", "CONTROL PLANE")
    add_image_crop(slide, ASSETS / "aihub-governance-navigation.png", 0.7, 1.55, 11.93, 3.5)
    add_rect(slide, 0.7, 1.55, 11.93, 3.5, fill=BG, radius=False, alpha=6)
    nav_groups = [
        ("BUILD", "安裝套件", "能力供應", BLUE),
        ("ACCESS", "客戶／API 金鑰", "誰能用", CYAN),
        ("CLUSTER", "節點管理", "在哪裡跑", TEAL),
        ("TEST", "測試中心", "能不能用", AMBER),
        ("TRACE", "記錄中心", "發生什麼", GREEN),
        ("CONFIG", "環境／設定", "怎麼維運", BLUE),
    ]
    x = 0.72
    for tag, title, body, accent in nav_groups:
        add_rect(slide, x, 5.43, 1.79, 0.92, fill=PANEL, line=accent)
        add_text(slide, tag, x + 0.14, 5.59, 0.58, 0.19, 8, accent, True, DISPLAY)
        add_text(slide, title, x + 0.14, 5.89, 1.46, 0.23, 11, TEXT, True)
        add_text(slide, body, x + 0.98, 5.6, 0.62, 0.22, 9, MUTED, True, align=PP_ALIGN.RIGHT)
        x += 1.99
    add_footer(slide, 23, "實機導覽截圖：控制台、Pack、客戶、Token、Cluster、測試、記錄、環境與設定。")

    # 24 Records
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 24, "記錄中心，是 AI 平台的黑盒子與學習回路", "OBSERVABILITY")
    log_sources = [
        ("執行歷程", "任務跑了什麼", BLUE),
        ("API 記錄", "誰呼叫哪個 Mode", CYAN),
        ("背景工作", "排程與長任務", TEAL),
        ("服務記錄", "容器與模型狀態", AMBER),
        ("系統記錄", "平台事件與錯誤", RED),
    ]
    y = 1.58
    for title, body, accent in log_sources:
        add_rect(slide, 0.76, y, 2.52, 0.79, fill=PANEL, line=accent)
        add_text(slide, title, 0.94, y + 0.15, 1.08, 0.25, 13, TEXT, True)
        add_text(slide, body, 1.93, y + 0.17, 1.08, 0.22, 9, MUTED, align=PP_ALIGN.RIGHT)
        add_line(slide, 3.28, y + 0.4, 5.02, 3.53, accent, 1.2)
        y += 0.94
    add_circle(slide, 5.02, 2.43, 2.2, fill=BG2, line=CYAN)
    add_text(slide, "TRACE", 5.02, 3.0, 2.2, 0.38, 22, CYAN, True, DISPLAY, PP_ALIGN.CENTER)
    add_text(slide, "可還原的事件鏈", 5.02, 3.48, 2.2, 0.3, 12, TEXT, True, align=PP_ALIGN.CENTER)
    decisions = [
        ("維運", "哪個服務壞了？", AMBER),
        ("容量", "哪台主機太忙？", CYAN),
        ("治理", "誰用了多少？", GREEN),
        ("產品", "哪個能力最需要？", TEAL),
    ]
    y = 1.72
    for title, body, accent in decisions:
        add_rect(slide, 8.4, y, 3.87, 0.92, fill=PANEL, line=accent)
        add_text(slide, title, 8.62, y + 0.19, 0.68, 0.27, 14, accent, True)
        add_text(slide, body, 9.37, y + 0.2, 2.55, 0.25, 13, TEXT, True)
        add_line(slide, 7.22, 3.53, 8.4, y + 0.46, accent, 1.2)
        y += 1.14
    add_text(slide, "沒有記錄，只能猜；有了記錄，才可以改善、計價與治理。", 4.08, 6.2, 8.2, 0.42, 17, GREEN, True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 24)

    # 25 Role: ordinary engineer
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 25, "普通顯卡工程師：不再被環境安裝困在第一公里", "ROLE 01")
    add_rect(slide, 0.72, 1.58, 5.72, 4.85, fill=RISK_PANEL, line=RED)
    add_text(slide, "BEFORE", 0.98, 1.85, 1.0, 0.28, 11, RED, True, DISPLAY)
    add_person(slide, 1.36, 2.45, 1.2, RED)
    add_rect(slide, 3.47, 2.25, 2.3, 1.42, fill=BG2, line=RED)
    add_text(slide, "CUDA？\n模型？\nDocker？", 3.47, 2.47, 2.3, 0.92, 18, TEXT, True, align=PP_ALIGN.CENTER)
    add_text(slide, "一半時間在裝環境，\n還沒做到功能就先放棄。", 1.14, 5.08, 4.9, 0.78, 19, MUTED, True)
    add_rect(slide, 6.86, 1.58, 5.75, 4.85, fill=SUCCESS_PANEL, line=GREEN)
    add_text(slide, "AFTER", 7.14, 1.85, 1.0, 0.28, 11, GREEN, True, DISPLAY)
    add_person(slide, 7.45, 2.45, 1.2, GREEN)
    add_rect(slide, 9.47, 2.28, 2.55, 1.35, fill=BG2, line=GREEN)
    add_text(slide, "curl / SDK\n→ AIHub API", 9.47, 2.64, 2.55, 0.64, 17, TEXT, True, font=DISPLAY, align=PP_ALIGN.CENTER)
    add_text(slide, "家裡的 GPU 仍可貢獻，\n重心回到應用與創意。", 7.28, 5.08, 4.9, 0.78, 19, GREEN, True)
    add_footer(slide, 25, "角色漫畫為情境示意；實際可用服務與效能依節點硬體與授權而定。")

    # 26 Role: cluster engineer
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 26, "Cluster 工程師：數十台主機，開始像一個產品運作", "ROLE 02")
    add_person(slide, 0.92, 2.28, 1.35, CYAN)
    add_text(slide, "從逐台 SSH…", 0.74, 5.09, 2.8, 0.36, 18, RED, True, align=PP_ALIGN.CENTER)
    add_chevron(slide, 3.3, 3.08, 0.55, 0.82, CYAN)
    add_circle(slide, 4.35, 2.25, 2.3, fill=BG2, line=CYAN)
    add_text(slide, "AIHub", 4.35, 2.88, 2.3, 0.4, 25, TEXT, True, DISPLAY, PP_ALIGN.CENTER)
    add_text(slide, "統一治理平面", 4.35, 3.4, 2.3, 0.3, 12, MUTED, True, align=PP_ALIGN.CENTER)
    node_positions = [(8.0, 1.78), (10.52, 1.78), (8.0, 4.15), (10.52, 4.15)]
    for i, (x, y) in enumerate(node_positions):
        accent = [CYAN, TEAL, BLUE, GREEN][i]
        add_rect(slide, x, y, 1.72, 1.38, fill=PANEL, line=accent)
        add_text(slide, f"NODE {i + 1:02d}", x + 0.18, y + 0.18, 1.2, 0.23, 9, accent, True, DISPLAY)
        add_text(slide, "GPU／Service\nHealth／VRAM", x + 0.18, y + 0.58, 1.26, 0.55, 11, TEXT, True)
        add_line(slide, 6.65, 3.4, x, y + 0.69, accent, 1.4)
    add_text(slide, "到看見健康、容量、服務與使用者。", 3.8, 5.84, 8.4, 0.42, 20, GREEN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 26)

    # 27 Role: PM and planner
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 27, "PM＋規劃師：知道地端能力，創意才有落點", "ROLE 03")
    add_person(slide, 0.9, 2.15, 1.3, AMBER)
    add_rect(slide, 3.0, 1.66, 3.04, 1.38, fill=PANEL, line=CYAN)
    add_text(slide, "「能不能把 3,000 份文件\n變成可查詢的知識？」", 3.25, 2.02, 2.52, 0.72, 17, TEXT, True)
    add_rect(slide, 3.0, 3.36, 3.04, 1.38, fill=PANEL, line=TEAL)
    add_text(slide, "「能不能把現勘影像\n自動轉成圖台線索？」", 3.25, 3.72, 2.52, 0.72, 17, TEXT, True)
    add_rect(slide, 3.0, 5.06, 3.04, 1.12, fill=PANEL, line=AMBER)
    add_text(slide, "「能不能先做 Demo？」", 3.25, 5.42, 2.52, 0.34, 17, TEXT, True)
    idea_flow = [
        ("問題", "需求語言", BLUE),
        ("Market", "能力組合", CYAN),
        ("Playground", "快速驗證", AMBER),
        ("API", "專案整合", TEAL),
        ("Usage", "量化價值", GREEN),
    ]
    y = 1.7
    for tag, body, accent in idea_flow:
        add_rect(slide, 7.04, y, 4.82, 0.77, fill=PANEL, line=accent)
        add_text(slide, tag, 7.28, y + 0.17, 1.1, 0.27, 14, accent, True, font=DISPLAY)
        add_text(slide, body, 8.58, y + 0.17, 2.9, 0.27, 15, TEXT, True)
        y += 0.98
    add_text(slide, "平台把「不知道能不能」變成「今天先跑一次」。", 7.04, 6.05, 5.28, 0.42, 18, GREEN, True)
    add_footer(slide, 27)

    # 28 Cyber runner
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 28, "「2077 竄網使」：Agent 在治理邊界內穿梭工具", "AGENT VISION")
    add_person(slide, 0.86, 2.32, 1.35, CYAN, fill=BG2)
    add_text(slide, "AGENT", 0.82, 5.1, 2.6, 0.3, 13, CYAN, True, DISPLAY, PP_ALIGN.CENTER)
    tool_nodes = [
        ("WEB", 4.0, 1.7, CYAN),
        ("OCR", 6.28, 1.7, TEAL),
        ("GEO", 8.56, 1.7, AMBER),
        ("VISION", 4.0, 4.28, BLUE),
        ("LLM", 6.28, 4.28, CYAN),
        ("TTS", 8.56, 4.28, GREEN),
    ]
    for tag, x, y, accent in tool_nodes:
        add_circle(slide, x, y, 1.25, fill=BG2, line=accent)
        add_text(slide, tag, x, y + 0.42, 1.25, 0.3, 13, accent, True, DISPLAY, PP_ALIGN.CENTER)
        add_line(slide, 2.62, 3.48, x, y + 0.63, accent, 1.2)
    add_rect(slide, 10.56, 2.2, 1.85, 2.35, fill=PANEL, line=GREEN)
    add_text(slide, "OUTPUT", 10.78, 2.5, 1.4, 0.26, 10, GREEN, True, DISPLAY)
    add_text(slide, "地圖線索\n任務摘要\n語音回報\n完整 Trace", 10.78, 3.0, 1.35, 1.15, 14, TEXT, True)
    add_text(slide, "Token、Allowed tools、核准與 Trace，讓 Agent 不是無限制亂竄。", 3.96, 6.15, 8.48, 0.36, 16, AMBER, True)
    add_footer(slide, 28, "願景示意：Agent 工具選擇仍受權限、核准與平台政策限制。")

    # 29 Hermes era
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 29, "從「小龍蝦」到 Hermes：工具會用，還要會協作", "AGENT ERA")
    add_rect(slide, 0.74, 1.58, 4.78, 4.86, fill=RISK_PANEL, line=RED)
    add_text(slide, "小龍蝦時代", 1.04, 1.94, 3.9, 0.42, 24, RED, True)
    add_text(slide, "一個指令\n一個工具\n一次結果", 1.04, 2.7, 3.9, 1.28, 25, TEXT, True)
    add_text(slide, "人手動換頁、貼資料、\n追錯誤、補上下文。", 1.04, 4.6, 3.72, 0.9, 17, MUTED, True)
    add_chevron(slide, 5.93, 3.22, 0.72, 1.05, CYAN)
    add_rect(slide, 6.99, 1.58, 5.62, 4.86, fill=SUCCESS_PANEL, line=GREEN)
    add_text(slide, "Hermes／Agent 時代", 7.31, 1.94, 4.8, 0.42, 24, GREEN, True)
    agent_traits = [
        ("PLAN", "拆解多步任務"),
        ("TOOLS", "選擇允許的工具"),
        ("HANDOFF", "交給合適角色"),
        ("TRACE", "留下完整行動軌跡"),
    ]
    y = 2.74
    for tag, body in agent_traits:
        add_text(slide, tag, 7.32, y, 0.9, 0.25, 9, CYAN, True, DISPLAY)
        add_text(slide, body, 8.32, y - 0.02, 3.34, 0.29, 15, TEXT, True)
        y += 0.68
    add_text(slide, "AIHub 提供的，是可治理、可替換的工具底座。", 7.32, 5.64, 4.62, 0.42, 17, GREEN, True)
    add_footer(slide, 29, "Hermes 為 Agent 願景示例；OpenAI Agents SDK 支援 tools、handoffs、guardrails 與 tracing。")

    # 30 Agent field workflow
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 30, "一個現勘任務，可以串成七段自動化服務", "AGENT DEMO")
    workflow = [
        ("01", "接收任務", "自然語言", BLUE),
        ("02", "開網頁", "Screenshot", CYAN),
        ("03", "讀內容", "OCR／Doc", TEAL),
        ("04", "洗地址", "Taiwan GEO", AMBER),
        ("05", "看影像", "Vision", BLUE),
        ("06", "產摘要", "LLM／RAG", CYAN),
        ("07", "回報", "TTS＋Trace", GREEN),
    ]
    x = 0.6
    for i, (num, title, body, accent) in enumerate(workflow):
        add_rect(slide, x, 2.04, 1.47, 2.68, fill=PANEL, line=accent)
        add_text(slide, num, x + 0.18, 2.26, 0.38, 0.23, 9, accent, True, DISPLAY)
        add_text(slide, title, x + 0.18, 2.88, 1.1, 0.64, 17, TEXT, True)
        add_text(slide, body, x + 0.18, 4.1, 1.12, 0.28, 10, accent, True, DISPLAY)
        if i < len(workflow) - 1:
            add_chevron(slide, x + 1.55, 3.1, 0.22, 0.44, accent)
        x += 1.79
    add_rect(slide, 1.08, 5.4, 11.2, 0.82, fill=BG2, line=GREEN)
    add_text(slide, "每一步都能替換模型、改節點、做核准；應用不用重寫整套基礎設施。", 1.08, 5.64, 11.2, 0.34, 17, GREEN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 30, "情境示意；Responses API 可使用 function calling 與 remote MCP 連接工具，敏感工具可要求核准。")

    # 31 Future market
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, 31, "未來 Market：能力越多，治理要比模型成長得更快", "FUTURE")
    add_text(slide, "能力層", 0.78, 1.56, 1.0, 0.28, 11, CYAN, True, DISPLAY)
    future_tools = [
        ("IMAGE", "圖片生成", CYAN),
        ("VIDEO", "影片生成", BLUE),
        ("3D", "3D 建模", TEAL),
        ("SPATIAL", "空間 AI", GREEN),
        ("DOC", "文件智慧", AMBER),
        ("AGENT", "多 Agent 協作", CYAN),
    ]
    x = 0.78
    for tag, title, accent in future_tools:
        add_circle(slide, x, 2.02, 1.35, fill=BG2, line=accent)
        add_text(slide, tag, x, 2.43, 1.35, 0.25, 10, accent, True, DISPLAY, PP_ALIGN.CENTER)
        add_text(slide, title, x - 0.15, 3.54, 1.65, 0.34, 15, TEXT, True, align=PP_ALIGN.CENTER)
        x += 2.03
    add_text(slide, "治理層", 0.78, 4.42, 1.0, 0.28, 11, GREEN, True, DISPLAY)
    governance = [
        ("REGISTRY", "版本／來源"),
        ("POLICY", "權限／核准"),
        ("EVALS", "品質／安全"),
        ("CAPACITY", "VRAM／排程"),
        ("ECONOMICS", "Quota／計價"),
    ]
    x = 0.78
    for tag, body in governance:
        add_rect(slide, x, 4.9, 2.25, 1.02, fill=PANEL, line=GREEN)
        add_text(slide, tag, x + 0.17, 5.1, 0.92, 0.23, 9, GREEN, True, DISPLAY)
        add_text(slide, body, x + 0.17, 5.5, 1.84, 0.25, 13, TEXT, True)
        x += 2.42
    add_text(slide, "新模型帶來驚喜；治理讓驚喜可被企業放心使用。", 0.78, 6.38, 11.8, 0.4, 18, GREEN, True, align=PP_ALIGN.CENTER)
    add_footer(slide, 31, "圖片生成、影片生成、3D 與多 Agent 協作為未來能力展望，非現況承諾。")

    # 32 Close
    slide = prs.slides.add_slide(blank)
    add_image_crop(slide, key_visual, 0, 0, W, H)
    add_rect(slide, 0, 0, W, H, fill=BG, radius=False, alpha=72)
    add_text(slide, "32  CALL TO ACTION", 0.72, 0.52, 3.0, 0.28, 11, CYAN, True, DISPLAY)
    add_text(slide, "下一台 AI 主機，不再是孤島", 0.72, 1.18, 8.3, 0.72, 35, TEXT, True)
    add_text(
        slide,
        "一台強的電腦，也走不了很遠。\n一群被治理、能協作的電腦，可以走得很遠。",
        0.74,
        2.24,
        7.7,
        1.36,
        26,
        TEXT,
        True,
    )
    ctas = [
        ("INTEGRATION", "把 GPU 連起來", BLUE),
        ("INTEROPERABILITY", "把服務標準化", CYAN),
        ("COMMUNICATION", "把狀態與用量看清楚", AMBER),
        ("SHARING", "把多餘算力共享出去", GREEN),
    ]
    x = 0.74
    for tag, body, accent in ctas:
        add_rect(slide, x, 4.25, 2.86, 1.02, fill=BG2, line=accent, alpha=94)
        add_text(slide, tag, x + 0.18, 4.46, 1.15, 0.22, 8, accent, True, DISPLAY)
        add_text(slide, body, x + 0.18, 4.82, 2.44, 0.28, 13, TEXT, True)
        x += 3.04
    add_text(slide, "買的是 GPU，累積的是企業能力。", 0.74, 5.88, 7.2, 0.56, 24, GREEN, True)
    add_text(slide, "AI 算了吧...  ·  GIDAY 08.07", 0.74, 6.58, 4.8, 0.32, 13, MUTED, True)
    add_text(slide, "32", 12.2, 7.02, 0.42, 0.24, 9, CYAN, True, DISPLAY, PP_ALIGN.RIGHT)

    assert len(prs.slides) == 32
    assert ("淡色版" in OUT.name) == LIGHT_THEME
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build_deck()
