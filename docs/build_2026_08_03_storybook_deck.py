#!/usr/bin/env python3
from __future__ import annotations

import math
import random
from functools import lru_cache
from pathlib import Path

from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.util import Inches


ROOT = Path('/var/www/html/3waAIHub/docs')
ASSETS = ROOT / '2026-08-03分享會投影片-assets'
OUT = ROOT / '2026-08-03分享會投影片.pptx'
MASTER = ROOT / '2026-08-03分享會投影片-主視覺.png'
W, H = 1920, 1080

PAPER = '#F1E2BF'
PAPER_LIGHT = '#FAF1D9'
INK = '#2A211C'
MUTED = '#6F6254'
BROWN = '#594034'
WOOD = '#34251F'
BLUE = '#1D83B5'
BLUE_LIGHT = '#8DD8EC'
ORANGE = '#D76532'
GOLD = '#B98C45'
GREEN = '#3E7A64'
RED = '#A9443A'
METAL = '#6E777A'

SERIF_REG = '/usr/share/fonts/opentype/noto/NotoSerifCJK-Regular.ttc'
SERIF_BOLD = '/usr/share/fonts/opentype/noto/NotoSerifCJK-Bold.ttc'
SANS_REG = '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc'
SANS_BOLD = '/usr/share/fonts/opentype/noto/NotoSansCJK-Bold.ttc'
MONO = '/home/john/.codex/skills/canvas-design/canvas-fonts/DMMono-Regular.ttf'


@lru_cache(maxsize=None)
def font(path: str, size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(path, size)


def s(page, chapter, title, main, scene, composition, narration, transition, source=''):
    return {
        'page': page,
        'chapter': chapter,
        'title': title,
        'main': main,
        'scene': scene,
        'composition': composition,
        'narration': narration,
        'transition': transition,
        'source': source,
    }


SLIDES = [
    s(1, '序幕', 'AI 冒險與電網', '從一個工程師的 AI 冒險，\n到一座 AI 能力電網', 'cover',
      '暮色舞台中央，一位旅行老工程師騎著載滿零件與巨書的腳踏車入場；孩子們在書頁光線前形成觀眾席。',
      '今天不先講架構，也不先講產品功能。我想從一位工程師反覆解決專案問題的故事開始，看一套工具怎麼慢慢長成組織可以共用的 AI 基礎設施。',
      '故事開始以前，先聽見遠處一聲腳踏車鈴。'),
    s(2, '第一章｜阿伯來了', '奇怪的腳踏車', '叮鈴——\n它不是來送貨的。', 'arrival',
      '老腳踏車從左側駛入，輪軸掛著網路線、硬碟、GPU 與工具箱；右側留下大片舞台空間。',
      '一台很奇怪的老腳踏車來到會場。車上沒有商品，只有工程師會認得的東西：網路線、硬碟、工具、顯示卡，還有修了很多次仍捨不得丟的零件。',
      '但孩子們最先看到的，不是那些零件。'),
    s(3, '第一章｜阿伯來了', '車後的大書', '比伺服器還重的一本書。', 'big_book',
      '腳踏車後架放著一冊帶金屬鉸鏈的巨書，書脊像機櫃，扣環連著發光網路線。',
      '腳踏車後面，載著一本大得不像話的書。它既像旅行箱，也像機櫃；書頁裡夾著模型、容器、API 和一堆只有工程師才會留下的失敗紀錄。',
      '當阿伯停好車，幾個孩子立刻圍了上來。'),
    s(4, '第一章｜阿伯來了', '孩子們圍上來', '工程師｜PM｜使用者｜管理者', 'kids',
      '四位孩子以不同道具區分角色：扳手、計畫板、手機與望遠鏡，圍著尚未打開的故事書。',
      '這群孩子代表今天在場的不同角色。有人關心能不能做，有人關心何時交付，有人只想把服務用好，也有人想知道設備與投資到底產生了什麼價值。',
      '其中一位孩子忍不住先問了。'),
    s(5, '第一章｜阿伯來了', '這是 AI 嗎？', '「這是 AI 嗎？」', 'question',
      '孩子的提問以立體紙泡泡浮在書上方，阿伯扶著車把，還沒有打開書。',
      '這是一個很自然的問題。因為車上有 GPU、有模型、有伺服器，大家很容易把整件事情理解成又一個 AI。',
      '阿伯搖搖頭，只回答一句。'),
    s(6, '第一章｜阿伯來了', '阿伯的答案', '這不是一個 AI。\n這是一個被專案一步一步逼出來的平台。', 'answer',
      '巨書第一次微微打開，書縫透出藍光；阿伯與孩子只剩剪影，兩句話成為舞台主體。',
      '3waAIHub 不是某一個模型，也不是突然想像出來的產品。它是在一個又一個專案裡，被重複安裝、重複失敗、重複交付的現實一步一步逼出來的。',
      '故事的第一位真正主角，叫做小羽。'),
    s(7, '第二章｜每案重來', '小羽接案了', '需求很簡單：\n「讓 AI 幫忙。」', 'project',
      '年輕工程師小羽坐在桌前，需求卷軸只有一行字，背後卻堆著未知的技術箱。',
      '每個專案的開頭常常很像：需求聽起來很簡單，可能是辨識圖片、讀文件、產生文字，或訓練一個模型。真正的複雜度，都藏在這句話後面。',
      '小羽第一件事，當然是開始找方法。'),
    s(8, '第二章｜每案重來', '先找方法', 'Paper、GitHub、Demo。', 'search',
      '紙頁化成搜尋桌，論文、GitHub 分支與 Demo 視窗像卡片沿著放大鏡軌道展開。',
      '他先找論文，再找有沒有現成專案，接著看 Demo 能不能接近需求。這一步很有創造力，也很像探險；但從「看起來能跑」到「真的能交付」，中間還很遠。',
      '找到模型以後，下一個動作通常是下載。'),
    s(9, '第二章｜每案重來', '模型下載中', '權重很大，\n希望也很大。', 'download',
      '巨大模型權重像沉重貨箱沿進度條落下，小羽在一旁扶著快倒下的桌面。',
      '模型權重可能幾百 MB，也可能幾十 GB。下載完成的那一刻很有希望，但我們還不知道它需要哪個 Python、哪個 CUDA、哪一組套件，以及哪一張顯示卡。',
      '答案往往藏在 README 裡。'),
    s(10, '第二章｜每案重來', 'README 迷宮', 'pip install\n只是入口。', 'maze',
      'README 被畫成多層紙迷宮，入口寫著 pip install，岔路標記 requirements、patch 與 issue。',
      'README 通常會告訴我們怎麼開始，卻不一定能描述每一種主機與每一個版本差異。小羽從 requirements 走到 issue，再從 issue 繞回某個三年前的 commit。',
      '接著，真正熟悉的環境問題登場了。'),
    s(11, '第二章｜每案重來', '環境第一關', 'Python 版本不對。', 'conda',
      '不同版本的 Python 方塊堆成不穩的塔，Conda 環境圈試圖把它們固定。',
      'Conda 很有用，它讓不同專案可以擁有自己的 Python 與套件。但只要上游套件、編譯器或模型版本不同，環境仍然可能在另一台機器上長得不一樣。',
      '即使 Python 對了，顯示卡還有自己的版本世界。'),
    s(12, '第二章｜每案重來', '顯卡第二關', 'CUDA：\n「我也有版本。」', 'cuda',
      'GPU 像一扇機械關卡，CUDA、Driver、Torch 三個齒輪沒有完全咬合，小羽拿著扳手校正。',
      'Driver、CUDA、PyTorch 和模型框架要彼此相容。只差一個版本，可能就從可以推論變成找不到 library，或在最需要它的時候直接爆掉。',
      '最後，小羽終於讓模型動了。'),
    s(13, '第二章｜每案重來', '三分鐘與三天', '模型只跑了三分鐘，\n環境裝了三天。', 'timer',
      '左邊小沙漏代表三分鐘推論，右邊巨大日曆翻了三天；比例差異形成幽默。',
      '這句話並不誇張。真正的推論可能很快，但為了得到那三分鐘結果，我們花了三天準備環境、修版本、重抓模型，還要記得自己改過什麼。',
      '更精彩的是，下一個專案很快就來了。'),
    s(14, '第二章｜每案重來', '下一個專案', '全部重來。', 'reset',
      '前一頁完成的環境像紙牌屋被翻頁風吹散，新需求卷軸又回到起點。',
      '下一個專案可能用另一個模型、另一套框架、另一張 GPU。前一次的知識留在某個人的筆記與 terminal history，團隊卻再次從搜尋與安裝開始。',
      '小羽開始想：能不能讓下一次少重來一點？'),
    s(15, '第三章｜環境到服務', '換一個方法', '從工具，走向服務。', 'evolution',
      '故事書升起三階紙台：環境、容器、服務，小羽沿階梯推著模型前進。',
      '第一個轉折不是做更大的平台，而是讓模型從「只有我會跑」變成「別人也能穩定使用」。這需要先把環境和啟動方式固定下來。',
      '我們先從環境本身開始縮小。'),
    s(16, '第三章｜環境到服務', '環境變小了', 'Anaconda → Miniconda', 'miniconda',
      '大型行李箱 Anaconda 收斂成精簡工具箱 Miniconda，只保留專案需要的零件。',
      '從完整 Anaconda 到 Miniconda，是把環境變得更精簡、更容易搬動。它不是最後答案，但開始讓專案願意明確說出自己真正需要哪些依賴。',
      '接著，我們把環境的內容寫下來。'),
    s(17, '第三章｜環境到服務', '開始可重現', '匯出環境，\n記住依賴。', 'manifest',
      '一張像航海圖的環境清單被釘在書頁上，版本與 hash 成為可重建路線。',
      'environment.yml、requirements.txt 或 lock file，讓環境不再只存在某台電腦。可重現性變好，但作業系統、系統套件與 GPU runtime 還在外面。',
      '於是，模型被放進更完整的盒子。'),
    s(18, '第三章｜環境到服務', '模型進盒子', 'Docker / Container', 'docker',
      '透明紙盒包住模型、Python 與系統依賴，外側印著 Docker 與固定 image tag。',
      'Docker 把程式、系統套件與啟動方式包在一起。相同 image 可以在符合條件的主機上重建相近環境，交付和回復都比手工安裝可靠。',
      '盒子裝好了，下一步是讓它能被呼叫。'),
    s(19, '第三章｜環境到服務', '盒子會說話', '固定啟動，\n提供 API。', 'api_box',
      '容器盒伸出標準化管線，端點標著 health 與 invoke，旁邊亮起 ready 燈。',
      '當容器有固定啟動方式、健康檢查與 API，它就從一段程式變成服務。其他系統不必知道裡面用哪個框架，只要遵守輸入與輸出契約。',
      '可是，服務一多，新的問題又出現了。'),
    s(20, '第三章｜環境到服務', '但盒子散落', '每台主機，\n各自一套。', 'scattered',
      '不同主機島上散落著 OCR、TTS、YOLO 等容器盒，線路彼此不相通。',
      'Docker 讓每個服務比較好部署，卻沒有自動告訴我們：服務在哪台主機、誰能使用、版本是否一致、現在健不健康，以及出了問題要去哪裡找。',
      '孩子看著滿地的盒子，開始連續發問。'),
    s(21, '第三章｜環境到服務', '還少了管理', 'Docker 解決了部署，\n但沒有解決管理。', 'management_gap',
      '四個問號圍住散落的盒子：誰能用、用了多少、壞了誰知、怎麼更新。',
      '我們需要的已經不是更多 Docker，而是一個管理 AI 服務的系統。它要知道服務、使用者、任務、成果和硬體狀態，並把這些事情放進同一套治理流程。',
      '阿伯翻到書中第一次出現名字的那一頁。'),
    s(22, '第四章｜AIHub 誕生', '翻到新章節', '管理 AI 服務的系統。', 'hub_reveal',
      '巨書中央升起一座尚未接線的調度站，只有 3waAIHub 字樣與幾個空插槽。',
      '3waAIHub 一開始做的事情很務實：把已經能跑的模型與服務，用共同規格裝進平台，再補上原本散落在人工流程裡的管理工作。',
      '第一步，還是從服務盒開始。'),
    s(23, '第四章｜AIHub 誕生', '先做服務盒', 'AI Service Pack', 'service_pack',
      '一個 Pack 被拆成模型、環境、健康檢查、API 契約與設定五層紙件。',
      'AI Service Pack 不只是 Dockerfile。它描述用途、版本、硬體門檻、安裝方式、服務設定、健康檢查與可用 Mode，讓平台知道如何管理這項能力。',
      '有了共同規格，還要有共同入口。'),
    s(24, '第四章｜AIHub 誕生', '再做共同入口', 'API Gateway', 'gateway',
      '中央城門接受不同請求，再把道路導向 OCR、Vision、ASR 與 TTS 工坊。',
      'API Gateway 讓應用面對的是穩定入口，而不是每一個容器的網址。服務可以換位置、更新版本，使用者不必重新理解整個機房。',
      '但共同入口不能等於所有人都能任意使用。'),
    s(25, '第四章｜AIHub 誕生', '每人一張票', '帳號 ＋ Token', 'token',
      '不同角色拿著帶有短碼的通行票，城門只辨識 Token，不暴露內部節點。',
      '帳號代表使用者與責任，Token 則讓程式安全地帶著身分呼叫服務。Token 可以撤銷、輪替與限制範圍，不必把內部服務直接攤在網路上。',
      '有身分之後，平台才能回答誰可以做什麼。'),
    s(26, '第四章｜AIHub 誕生', '知道誰能用', '權限 ＋ 使用量', 'usage',
      '通行票經過權限閘門，旁邊的機械表記錄 Mode、次數、成功與失敗。',
      '治理不是把門鎖死，而是把權限說清楚。哪個 Token 可以使用哪個 Mode、呼叫了幾次、成功還是失敗，都成為可查詢的紀錄。',
      '有些請求很快，有些工作卻不能立即完成。'),
    s(27, '第四章｜AIHub 誕生', '工作會排隊', 'Queue ＋ Task', 'queue',
      '任務卡排進候車月台，Task ID 像車票號碼；不同 Worker 在另一端領取。',
      '長工作不能綁住一個網頁請求。平台先回傳 Task ID，任務進入 Queue，Worker 再依資源與順序執行；使用者可以查狀態，而不是一直等著連線。',
      '任務完成後，還要把結果可靠地交回去。'),
    s(28, '第四章｜AIHub 誕生', '成果有去處', 'Callback ＋ Artifact\n＋ Retention', 'artifact',
      'Worker 把 WAV、圖片、模型與 JSON 成果送進標記貨架，Callback 信差把完成消息送回。',
      'Artifact 是任務真正產生的檔案；Callback 負責通知；Retention 則決定成果、來源與工作區何時保留、何時清理。這些規則讓長任務不只是跑完，而是能被追查與回收。',
      '把這些零件放在一起，我們才能說出平台的定位。'),
    s(29, '第四章｜AIHub 誕生', '它不是模型', '3waAIHub 不是一個模型。\n它管理服務、任務與算力。', 'hub_core',
      '中央 Hub 外圍依序環繞 Service、Task、Account、Token、Usage 與 Health 六個齒輪。',
      '模型是能力來源，3waAIHub 是治理與供應能力的系統。它不取代模型，而是讓不同模型用一致方法被安裝、授權、執行、觀測與更新。',
      '如果要找一個更容易記住的比喻，它像什麼？'),
    s(30, '第四章｜AIHub 誕生', 'AI 能力電網', '它更像是一座 AI 電網。', 'grid',
      'Vision、OCR、LLM、ASR、TTS、YOLO、SAM、去背與文件解析化為小型發電工坊，中央變電所將能力送往使用者。',
      '每個模型像不同類型的發電機，適合不同工作；AIHub 像輸電網與變電所，負責接入、轉送、權限、計量與狀態。不是每個人都要理解發電機，才能使用電。',
      '孩子看懂電網，卻又想到另一個問題。'),
    s(31, '第五章｜真正做事', '孩子又問了', '「AI 不就是\n問問題、拿答案嗎？」', 'chat_question',
      '孩子拿著聊天泡泡，書中另一側卻升起訓練機、文件工坊與語音工作台。',
      '聊天是重要的 AI 使用方式，但不是全部。很多專案真正需要的是處理一批資料、訓練模型、轉換格式，或產生可以下載與再利用的成果。',
      '而這些工作，通常不會在幾秒內結束。'),
    s(32, '第五章｜真正做事', '有些工作很久', '訓練｜解析｜生成｜轉換', 'long_jobs',
      '四座紙雕工坊同時運作：模型訓練、文件解析、語音生成與格式轉換，齒輪轉速各異。',
      '訓練可能幾十分鐘甚至幾小時，大型文件解析需要切頁與辨識，語音和影像生成也可能分段完成。這些工作需要真正的執行環境，而不只是把訊息轉給另一支 API。',
      '所以不能要求使用者一直守在網頁前面。'),
    s(33, '第五章｜真正做事', '網頁不能等', '工作可以久，\n畫面不能卡。', 'browser_wait',
      '瀏覽器前的孩子離開座位，任務卡仍在後台工坊推進，進度表持續更新。',
      '正確做法是先接受任務、回傳識別碼，再由背景執行。即使瀏覽器關閉，任務仍有自己的狀態、期限與成果；回來時可以繼續查詢。',
      '平台因此學會一套新的動作。'),
    s(34, '第五章｜真正做事', '平台學會排隊', '提交 → 排隊 → 執行 → 回報', 'queue_flow',
      '四段紙軌道從 Submit 經 Queue、Worker 到 Result，Task ID 沿途不變。',
      '使用者提交後得到 Task ID；Queue 保存順序與條件；Worker 執行；完成後回報狀態並登記 Artifact。失敗也會留下明確階段，而不是只得到一個模糊錯誤。',
      '要讓 Worker 安全工作，還需要幾個重要機制。'),
    s(35, '第五章｜真正做事', '工匠開始工作', 'Worker ＋ Lease\n＋ Heartbeat', 'taskflow',
      '第一張技術流程圖以故事工坊呈現：Queue 發任務、Worker 取得 Lease、Heartbeat 維持所有權、Timeout／Cancel 進入 Recovery。',
      'Lease 代表某個 Worker 在限定時間內擁有這個工作；Heartbeat 證明它仍活著。Timeout、Cancel 與 Recovery 防止同一任務被重複執行，也避免 GPU 因殘留程序被錯誤釋放。',
      '這套流程讓平台能承接真正的模型工作。'),
    s(36, '第五章｜真正做事', 'YOLO 的一天', 'Predict｜Train｜ONNX Export', 'yolo',
      '同一個 YOLO 工坊分出推論、訓練與模型匯出三條生產線，成果回到 Artifact 貨架。',
      '以 YOLO 為例，平台不只呼叫推論 API，也可以提交 Predict、Train、ONNX Export 等工作。每個任務有自己的資源、輸入、進度、成果與清理規則。',
      '所以，AIHub 的角色不能只被理解成 API 轉接。'),
    s(37, '第五章｜真正做事', '真正承接工作', '不只提供 API，\n還要真正承接工作。', 'workers',
      '多位紙雕工匠在 GPU、CPU、文件與影音工作台間分工，中央平台只負責調度與紀錄。',
      'Local Service 適合長駐型能力；Local Execution 適合被平台派發的工作。兩者都透過相同帳號、任務、紀錄與成果治理，讓應用不用自己重做基礎設施。',
      '當服務與任務都增加，一台主機開始承受壓力。'),
    s(38, '第六章｜一台不夠', '小發電機冒煙', 'VRAM 不夠了。', 'overheat',
      '單一 GPU 發電機同時拉著 LLM、Vision 與語音工坊，溫度表升高，橘紅紙煙冒出。',
      '模型越來越多，不同服務同時常駐，VRAM 很快就成為瓶頸。單機上可能同時存在等待、閒置與滿載，光看顯示卡型號無法理解真正容量。',
      '孩子提出最直覺的解法。'),
    s(39, '第六章｜一台不夠', '換更大的嗎？', '「換一台更大的？」\n「可以。」', 'bigger',
      '孩子指向一台更大的 GPU 發電機，阿伯簡短點頭，但遠處還有更大的模型影子。',
      '更大的主機當然有價值，很多工作確實需要更大的 VRAM。但如果治理方式不變，所有服務仍堆在同一台機器，容量只會晚一點再次用完。',
      '下一頁，模型的影子繼續長大。'),
    s(40, '第六章｜一台不夠', '下次仍不夠', '真正答案，\n不是永遠換大主機。', 'limits',
      '巨大模型箱再次壓滿新主機，旁邊其他主機卻仍留有空位，形成資源不均。',
      '不同工作需要不同硬體：Vision、語音、訓練、CPU 工具與文件處理不必擠在同一台。真正可持續的方向，是讓多台主機保留各自專長，又能被同一平台調度。',
      '於是，故事書第一次攤開三座機房。'),
    s(41, '第六章｜一台不夠', '讓主機一起做', '8 GB｜32 GB｜16 GB\n不是合併 VRAM，是分工。', 'three_nodes',
      '三座紙雕機房標示 GTX 1080 8 GB、RTX 5090 32 GB、RTX 5060 Ti 16 GB，以不同能力線接向中央入口。',
      '目前三個節點已持續回報：8 GB、32 GB 與 16 GB。Cluster 並不是把三張卡的 VRAM 直接相加，而是讓每台主機提供適合自己的服務與 Worker，再由入口做選擇。',
      '分工成立以後，平台還要知道怎麼選路。',
      '3waAIHub 本機 Cluster 快照，2026-08-03 14:57。'),
    s(42, '第六章｜一台不夠', 'Cluster 不是拼卡', '能力註冊＋任務路由\n＋健康狀態', 'cluster',
      '節點各自掛上 Vision、LLM、Voice、Train 能力牌，Router 依 Mode、健康與容量導流。',
      'Cluster 的核心是能力註冊、統一入口、任務路由、健康狀態與資源利用。服務在哪台主機是平台的責任；使用者只選自己有權限的能力。故障時也要留下狀態，而不是假裝所有節點都一樣。',
      '對外看起來，整個網路仍只有一扇門。'),
    s(43, '第六章｜一台不夠', '入口只有一個', '對使用者而言，入口只有一個。\n入口後方，可以是一整座 AI 機房。', 'machine_room',
      '前景是一座單一 API 城門，後方展開多排主機、Worker 與能力工坊，路線在門內分流。',
      '使用者不必知道模型在哪台主機，也不必持有每個節點的憑證。統一入口保留身分、權限與使用量；後方的 Router 才處理節點選擇與結果回傳。',
      '做到這裡，故事已經不只屬於一位工程師。'),
    s(44, '第七章｜走向治理', '一開始', '我們只是想讓模型跑起來。', 'islands',
      '小羽站在一座小島上點亮第一個模型；遠方各部門島嶼也有自己的設備與專案。',
      '每個部門都有自己的需求、設備與專業，這些投入都很重要。但如果能力只留在個人電腦或單一專案裡，下一個團隊仍看不到，也無法重複使用。',
      '後來，我們開始把能力帶到共同的交換所。'),
    s(45, '第七章｜走向治理', '後來', '讓能力被重複使用。\nAI 能力交換所', 'exchange',
      '不同部門把模型與服務帶到中央交換站，仍保留自己的設備，卻能發布能力與取得通行權。',
      'AI 能力交換所不是把所有設備收走，而是讓模型可以註冊、服務可以發布、部門可以保留專長，使用者則從共同入口取得被授權的能力。',
      '當能力開始流動，治理就成為共同語言。'),
    s(46, '第七章｜走向治理', '治理全景', '知道能力、權限、用量、\n任務、成果與健康。', 'architecture',
      '完整紙雕架構首次出現：人與 Agent 經帳號／Token 進入 Gateway；分流至 Service、Local Execution 與 Cluster；底層由 Queue、Worker、Artifact、Usage、Logs、Health 支撐。',
      '治理不是限制使用者，而是回答六個問題：有哪些能力、誰可以用、用了多少、任務是否成功、成果在哪裡、服務與設備是否健康。未來 Agent 也能把這些能力當成受控工具，但仍經過權限、核准與追蹤。',
      '把全景收回故事，我們只剩下一個真正的下一步。',
      'OpenAI Agents SDK 官方文件：tools、handoffs、guardrails、approvals、tracing；https://developers.openai.com/api/docs/guides/agents#compare-the-responses-api-and-agents-sdk（2026-08-03）。'),
    s(47, '第七章｜走向治理', '下一步', '讓整個組織的 AI 能力連起來。', 'city',
      '整本書完全展開，模型、Docker、API、GPU、主機、Token、任務與部門連成一座發光但克制的 AI 城市電網。',
      '這條路徑很簡單：從工具走向服務，從服務走向平台，從平台走向治理。下一步不是要求每個人都重新蓋一套，而是讓組織既有的主機、模型與經驗連起來。',
      '阿伯準備闔上書，留下最後一句話。'),
    s(48, '終場｜Q&A', '故事講完了。', '現在，換你問問題。', 'qa',
      '阿伯闔上巨書，腳踏車停在一旁；一位孩子舉手，書頁邊緣仍保留整座城市的微光。',
      '故事講完了。現在，換你問問題。3waAIHub 是從實戰問題長出來的 AI 數位治理服務平台。不是每個人都要蓋一座發電廠，但每個人都應該能使用電。',
      'Q&A。'),
]


def paper_canvas(seed: int) -> Image.Image:
    random.seed(seed)
    noise = Image.effect_noise((W, H), 12).convert('L')
    texture = ImageOps.colorize(noise, '#D5C29B', '#FFF8E7').convert('RGBA')
    base = Image.blend(Image.new('RGBA', (W, H), PAPER_LIGHT), texture, 0.16)
    d = ImageDraw.Draw(base, 'RGBA')
    d.rectangle((0, 0, W, H), fill=(48, 34, 27, 255))
    # Open book with a shallow perspective and physical page edges.
    d.polygon([(48, 70), (946, 38), (958, 1028), (60, 1000)], fill=PAPER, outline='#A88961')
    d.polygon([(974, 38), (1872, 70), (1860, 1000), (962, 1028)], fill=PAPER_LIGHT, outline='#A88961')
    for i in range(7):
        shade = 52 - i * 5
        d.line((955 + i, 56, 956 + i, 1016), fill=(57, 42, 35, shade), width=2)
    d.line((960, 52, 960, 1018), fill=(92, 65, 48, 105), width=5)
    # Travel-theatre frame and hardware.
    d.rounded_rectangle((28, 26, 1892, 1052), radius=28, outline=(184, 140, 73, 190), width=5)
    for x, y in [(45, 45), (1875, 45), (45, 1034), (1875, 1034)]:
        d.ellipse((x - 9, y - 9, x + 9, y + 9), fill=(108, 91, 72, 220), outline=(229, 193, 121, 220), width=2)
    for _ in range(140):
        x = random.randint(65, 1855)
        y = random.randint(55, 1018)
        r = random.choice([1, 1, 1, 2])
        d.ellipse((x-r, y-r, x+r, y+r), fill=(97, 75, 52, random.randint(18, 38)))
    return base


def rr(draw, box, radius=24, fill=PAPER_LIGHT, outline=BROWN, width=3, shadow=None):
    if shadow is not None:
        x1, y1, x2, y2 = box
        shadow.rounded_rectangle((x1 + 10, y1 + 14, x2 + 10, y2 + 14), radius=radius, fill=(36, 25, 19, 90))
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def line(draw, pts, fill=INK, width=5):
    draw.line(pts, fill=fill, width=width, joint='curve')


def label(draw, xy, text, size=24, color=INK, bold=False, mono=False, anchor='la'):
    path = MONO if mono else (SANS_BOLD if bold else SANS_REG)
    draw.text(xy, text, font=font(path, size), fill=color, anchor=anchor)


def person(draw, x, y, scale=1.0, color=BROWN, child=False, tool=None):
    head = 32 if child else 42
    draw.ellipse((x-head*scale, y-head*scale, x+head*scale, y+head*scale), fill=PAPER_LIGHT, outline=color, width=max(2, int(4*scale)))
    body_h = (118 if child else 150) * scale
    body_w = (64 if child else 78) * scale
    draw.rounded_rectangle((x-body_w, y+44*scale, x+body_w, y+44*scale+body_h), radius=int(22*scale), fill='#D8C49D', outline=color, width=max(2, int(4*scale)))
    line(draw, [(x-body_w, y+80*scale), (x-110*scale, y+135*scale)], color, max(2, int(5*scale)))
    line(draw, [(x+body_w, y+80*scale), (x+110*scale, y+135*scale)], color, max(2, int(5*scale)))
    line(draw, [(x-35*scale, y+44*scale+body_h), (x-65*scale, y+44*scale+body_h+90*scale)], color, max(2, int(5*scale)))
    line(draw, [(x+35*scale, y+44*scale+body_h), (x+65*scale, y+44*scale+body_h+90*scale)], color, max(2, int(5*scale)))
    if tool:
        label(draw, (x, y+105*scale), tool, int(22*scale), color, True, anchor='mm')


def storyteller(draw, x, y, scale=1.0):
    person(draw, x, y, scale, BROWN)
    draw.polygon([(x-70*scale, y-46*scale), (x+70*scale, y-46*scale), (x+34*scale, y-92*scale), (x-20*scale, y-88*scale)], fill='#6D4D34', outline=INK)
    draw.arc((x-42*scale, y-4*scale, x+42*scale, y+82*scale), 10, 170, fill='#B7A17D', width=max(3, int(7*scale)))
    line(draw, [(x+105*scale, y+135*scale), (x+142*scale, y+315*scale)], GOLD, max(3, int(6*scale)))


def bicycle(draw, x, y, scale=1.0, loaded=True):
    r = 90*scale
    for cx in (x, x+280*scale):
        draw.ellipse((cx-r, y-r, cx+r, y+r), outline=BROWN, width=max(3, int(7*scale)))
        for a in range(0, 360, 30):
            ex = cx + math.cos(math.radians(a))*r
            ey = y + math.sin(math.radians(a))*r
            line(draw, [(cx, y), (ex, ey)], '#927963', max(1, int(2*scale)))
    line(draw, [(x, y), (x+105*scale, y-95*scale), (x+205*scale, y), (x, y), (x+150*scale, y)], ORANGE, max(4, int(10*scale)))
    line(draw, [(x+105*scale, y-95*scale), (x+250*scale, y-130*scale)], BROWN, max(3, int(7*scale)))
    line(draw, [(x+205*scale, y), (x+280*scale, y)], BROWN, max(3, int(7*scale)))
    if loaded:
        rr(draw, (x-35*scale, y-260*scale, x+125*scale, y-135*scale), int(12*scale), '#6E5844', GOLD, max(2, int(4*scale)))
        label(draw, (x+45*scale, y-200*scale), 'GPU', int(24*scale), PAPER_LIGHT, True, True, 'mm')
        line(draw, [(x+10*scale, y-145*scale), (x-60*scale, y-95*scale)], BLUE, max(2, int(5*scale)))


def server(draw, x, y, w=180, h=260, accent=BLUE, title='NODE'):
    rr(draw, (x, y, x+w, y+h), 18, '#5C615F', INK, 4)
    for i in range(5):
        yy = y+40+i*38
        draw.rounded_rectangle((x+22, yy, x+w-22, yy+24), radius=6, fill='#272B2B', outline=accent, width=2)
        draw.ellipse((x+w-44, yy+7, x+w-34, yy+17), fill=accent)
    label(draw, (x+w/2, y+h+24), title, 20, INK, True, True, 'ma')


def model_box(draw, x, y, w=210, h=135, title='MODEL', accent=BLUE):
    rr(draw, (x, y, x+w, y+h), 18, '#E5D3AC', accent, 4)
    draw.polygon([(x+16, y+22), (x+w-16, y+12), (x+w-34, y+54), (x+28, y+60)], fill='#F7EBCF', outline=accent)
    label(draw, (x+w/2, y+h-38), title, 22, INK, True, True, 'mm')


def arrow(draw, a, b, color=BLUE, width=7):
    line(draw, [a, b], color, width)
    ang = math.atan2(b[1]-a[1], b[0]-a[0])
    head = 18
    p1 = (b[0]-head*math.cos(ang-math.pi/6), b[1]-head*math.sin(ang-math.pi/6))
    p2 = (b[0]-head*math.cos(ang+math.pi/6), b[1]-head*math.sin(ang+math.pi/6))
    draw.polygon([b, p1, p2], fill=color)


def bubble(draw, box, text, accent=BLUE, size=44):
    rr(draw, box, 34, '#FFF8E6', accent, 5)
    x1, y1, x2, y2 = box
    draw.polygon([(x1+95, y2), (x1+145, y2), (x1+110, y2+48)], fill='#FFF8E6', outline=accent)
    draw.multiline_text(((x1+x2)/2, (y1+y2)/2), text, font=font(SERIF_BOLD, size), fill=INK, anchor='mm', align='center', spacing=10)


def workshop(draw, x, y, title, accent=BLUE, scale=1.0):
    w, h = 190*scale, 150*scale
    rr(draw, (x, y, x+w, y+h), int(16*scale), '#DDC99F', accent, max(2, int(4*scale)))
    draw.polygon([(x-8*scale, y+28*scale), (x+w/2, y-34*scale), (x+w+8*scale, y+28*scale)], fill='#8E6A4D', outline=INK)
    draw.ellipse((x+w*0.38, y+45*scale, x+w*0.62, y+78*scale), outline=accent, width=max(2, int(4*scale)))
    label(draw, (x+w/2, y+h-30*scale), title, int(18*scale), INK, True, True, 'mm')


def draw_scene(draw, scene, box, page):
    x1, y1, x2, y2 = box
    cx, cy = (x1+x2)//2, (y1+y2)//2
    if scene in {'cover', 'arrival'}:
        bicycle(draw, cx-270, cy+190, 1.15, True)
        storyteller(draw, cx+235, cy-40, 1.0)
        for i in range(3):
            person(draw, x2-350+i*115, y2-230+(i%2)*18, .62, [BLUE, ORANGE, GREEN][i], True)
        if scene == 'cover':
            draw.arc((cx-180, cy-250, cx+430, cy+330), 205, 340, fill=BLUE_LIGHT, width=8)
    elif scene == 'big_book':
        bicycle(draw, cx-330, cy+220, 1.2, False)
        rr(draw, (cx-170, cy-260, cx+360, cy+180), 24, '#7C5738', GOLD, 8)
        for i in range(7):
            draw.line((cx-130, cy-210+i*48, cx+320, cy-210+i*48), fill='#D7BE8C', width=4)
        label(draw, (cx+95, cy-10), 'STORY / SYSTEM', 26, PAPER_LIGHT, True, True, 'mm')
        line(draw, [(cx-180, cy+70), (cx-300, cy+130)], BLUE, 8)
    elif scene == 'kids':
        roles = [('DEV', BLUE), ('PM', ORANGE), ('USER', GREEN), ('OPS', RED)]
        for i, (t, c) in enumerate(roles):
            px = x1+150+i*185
            person(draw, px, cy-80+(i%2)*24, .78, c, True, t)
        rr(draw, (cx-85, cy-250, cx+240, cy+50), 24, '#7C5738', GOLD, 7)
        label(draw, (cx+78, cy-100), 'THE BOOK', 24, PAPER_LIGHT, True, True, 'mm')
    elif scene in {'question', 'chat_question'}:
        person(draw, x1+180, cy+60, .85, BLUE, True, '?')
        storyteller(draw, x2-210, cy-20, .82)
        bubble(draw, (x1+320, y1+65, x2-300, y1+310), '這是 AI 嗎？' if scene == 'question' else 'AI 不就是\n問問題、拿答案嗎？', BLUE, 42)
    elif scene == 'answer':
        rr(draw, (cx-330, cy-260, cx+330, cy+270), 30, '#7C5738', GOLD, 7)
        line(draw, [(cx, cy-235), (cx, cy+245)], '#D6B778', 6)
        for i in range(4):
            draw.arc((cx-280+i*8, cy-210+i*8, cx+280-i*8, cy+230-i*8), 210, 330, fill=(29,131,181,100), width=3)
    elif scene == 'project':
        person(draw, x1+190, cy+20, .92, BLUE, False, '小羽')
        rr(draw, (x1+380, cy-240, x2-70, cy+120), 22, '#F8EBCF', GOLD, 5)
        label(draw, (x1+425, cy-175), 'PROJECT BRIEF', 22, ORANGE, True, True)
        label(draw, (x1+425, cy-90), '讓 AI 幫忙', 48, INK, True)
        for i in range(4):
            model_box(draw, x2-320+i*24, cy+175-i*12, 190, 105, ['?', 'DATA', 'GPU', 'TIME'][i], [RED, BLUE, ORANGE, GREEN][i])
    elif scene == 'search':
        draw.ellipse((cx-130, cy-130, cx+130, cy+130), outline=BLUE, width=14)
        line(draw, [(cx+90, cy+95), (cx+250, cy+250)], BLUE, 18)
        for i, (t, c, a) in enumerate([('PAPER', GOLD, -130), ('GITHUB', INK, 0), ('DEMO', ORANGE, 130)]):
            rr(draw, (cx-360+i*250, cy-290+a//4, cx-160+i*250, cy-145+a//4), 16, '#F6E8C8', c, 4)
            label(draw, (cx-260+i*250, cy-218+a//4), t, 24, c, True, True, 'mm')
    elif scene == 'download':
        model_box(draw, cx-180, cy-230, 360, 260, 'MODEL WEIGHTS', BLUE)
        rr(draw, (cx-380, cy+130, cx+380, cy+185), 20, '#D4C19D', BROWN, 3)
        draw.rounded_rectangle((cx-370, cy+140, cx+250, cy+175), radius=14, fill=BLUE)
        label(draw, (cx, cy+240), '18.6 GB / 20.0 GB', 24, MUTED, True, True, 'mm')
        person(draw, x2-180, cy+35, .7, ORANGE, False, '小羽')
    elif scene == 'maze':
        for i in range(7):
            yy = y1+90+i*90
            line(draw, [(x1+80+(i%2)*120, yy), (x2-80-(i%3)*90, yy)], BROWN, 8)
            if i < 6:
                xx = x1+180+(i*137)%(x2-x1-360)
                line(draw, [(xx, yy), (xx, yy+90)], BROWN, 8)
        rr(draw, (x1+80, y1+40, x1+310, y1+120), 16, '#FFF8E6', BLUE, 4)
        label(draw, (x1+195, y1+80), 'pip install', 24, BLUE, True, True, 'mm')
        label(draw, (x2-150, y2-70), 'ISSUE #427', 20, RED, True, True, 'mm')
    elif scene == 'conda':
        for i, t in enumerate(['PY 3.8', 'PY 3.10', 'PY 3.12']):
            model_box(draw, cx-310+i*205, cy-80-i*55, 180, 110, t, [GREEN, BLUE, RED][i])
        draw.arc((cx-390, cy-300, cx+390, cy+330), 10, 330, fill=GREEN, width=12)
        label(draw, (cx, cy+300), 'CONDA ENV', 28, GREEN, True, True, 'mm')
    elif scene == 'cuda':
        server(draw, cx-105, cy-200, 210, 300, ORANGE, 'GPU')
        for i, (t, c) in enumerate([('DRIVER', METAL), ('CUDA', ORANGE), ('TORCH', BLUE)]):
            xx = cx-290+i*290
            draw.ellipse((xx-65, cy+160-65, xx+65, cy+160+65), fill='#D7C59F', outline=c, width=7)
            label(draw, (xx, cy+160), t, 18, c, True, True, 'mm')
            if i < 2:
                line(draw, [(xx+65, cy+160), (xx+225, cy+160)], RED, 5)
    elif scene == 'timer':
        draw.ellipse((cx-370, cy-210, cx-70, cy+90), fill='#F6E9CC', outline=BLUE, width=8)
        line(draw, [(cx-220, cy-60), (cx-220, cy-155)], INK, 8)
        line(draw, [(cx-220, cy-60), (cx-145, cy-10)], ORANGE, 8)
        label(draw, (cx-220, cy+160), '3 MIN', 30, BLUE, True, True, 'mm')
        rr(draw, (cx+30, cy-230, cx+390, cy+110), 24, '#F4E5C4', RED, 7)
        for i in range(3):
            label(draw, (cx+210, cy-135+i*90), f'DAY {i+1}', 34, RED if i == 2 else BROWN, True, True, 'mm')
    elif scene == 'reset':
        draw.arc((cx-260, cy-250, cx+260, cy+270), 35, 330, fill=RED, width=16)
        draw.polygon([(cx+238, cy-70), (cx+300, cy-20), (cx+225, cy+20)], fill=RED)
        for i, t in enumerate(['ENV', 'CUDA', 'MODEL']):
            model_box(draw, cx-280+i*190, cy-80+(i%2)*70, 165, 100, t, [BLUE, ORANGE, GREEN][i])
        label(draw, (cx, cy+300), 'BACK TO ZERO', 26, RED, True, True, 'mm')
    elif scene == 'evolution':
        items = [('ENV', BROWN), ('CONTAINER', ORANGE), ('SERVICE', BLUE)]
        for i, (t, c) in enumerate(items):
            xx = x1+100+i*280
            model_box(draw, xx, cy-90-i*80, 225, 135, t, c)
            if i < 2:
                arrow(draw, (xx+235, cy-25-i*80), (xx+270, cy-25-(i+1)*80), c, 7)
    elif scene == 'miniconda':
        rr(draw, (x1+100, cy-190, x1+420, cy+170), 34, '#D6BD91', BROWN, 6)
        label(draw, (x1+260, cy-10), 'ANACONDA', 28, BROWN, True, True, 'mm')
        arrow(draw, (x1+450, cy), (x2-400, cy), BLUE, 10)
        rr(draw, (x2-360, cy-110, x2-110, cy+110), 28, '#F5E9CC', BLUE, 6)
        label(draw, (x2-235, cy), 'MINICONDA', 24, BLUE, True, True, 'mm')
    elif scene == 'manifest':
        rr(draw, (cx-350, cy-270, cx+350, cy+270), 22, '#F8EBCD', GOLD, 5)
        label(draw, (cx-290, cy-210), 'ENVIRONMENT MANIFEST', 24, BLUE, True, True)
        for i, text in enumerate(['python = 3.10', 'torch = pinned', 'cuda = compatible', 'sha256 = verified']):
            draw.rectangle((cx-290, cy-130+i*72, cx-250, cy-90+i*72), outline=GREEN, width=4)
            draw.line((cx-283, cy-108+i*72, cx-269, cy-94+i*72, cx-250, cy-123+i*72), fill=GREEN, width=5)
            label(draw, (cx-220, cy-110+i*72), text, 23, INK, False, True, 'lm')
    elif scene in {'docker', 'api_box'}:
        rr(draw, (cx-260, cy-220, cx+260, cy+180), 28, '#DCE7E8', BLUE, 8)
        model_box(draw, cx-125, cy-90, 250, 145, 'MODEL', ORANGE)
        label(draw, (cx, cy-175), 'CONTAINER', 24, BLUE, True, True, 'mm')
        if scene == 'api_box':
            for i, t in enumerate(['/health', '/invoke']):
                yy = cy-80+i*110
                arrow(draw, (cx+270, yy), (x2-80, yy), GREEN if i == 0 else BLUE, 8)
                label(draw, (x2-85, yy-34), t, 22, GREEN if i == 0 else BLUE, True, True, 'ra')
    elif scene in {'scattered', 'management_gap'}:
        for i, (t, c) in enumerate([('OCR', BLUE), ('TTS', ORANGE), ('YOLO', GREEN), ('LLM', RED), ('DOC', GOLD), ('SAM', METAL)]):
            xx = x1+60+(i%3)*255
            yy = y1+65+(i//3)*275
            model_box(draw, xx, yy, 205, 125, t, c)
            label(draw, (xx+103, yy+165), f'HOST {i+1}', 18, MUTED, True, True, 'mm')
        if scene == 'management_gap':
            for i, q in enumerate(['誰能用？', '用了多少？', '壞了誰知？', '怎麼更新？']):
                label(draw, (x1+120+(i%2)*390, y1+220+(i//2)*290), q, 28, RED, True, anchor='mm')
    elif scene == 'hub_reveal':
        draw.ellipse((cx-220, cy-220, cx+220, cy+220), fill='#D6E4E2', outline=BLUE, width=10)
        rr(draw, (cx-130, cy-115, cx+130, cy+115), 26, '#56605F', BLUE_LIGHT, 6)
        label(draw, (cx, cy), '3waAIHub', 32, PAPER_LIGHT, True, True, 'mm')
        for a in range(0, 360, 45):
            ex = cx+math.cos(math.radians(a))*330
            ey = cy+math.sin(math.radians(a))*250
            line(draw, [(cx, cy), (ex, ey)], BLUE_LIGHT, 4)
            draw.ellipse((ex-24, ey-24, ex+24, ey+24), fill=PAPER_LIGHT, outline=BLUE, width=4)
    elif scene == 'service_pack':
        parts = [('MODEL', ORANGE), ('RUNTIME', METAL), ('HEALTH', GREEN), ('API', BLUE), ('SETTINGS', GOLD)]
        for i, (t, c) in enumerate(parts):
            xx = x1+50+i*145
            yy = cy+100-i*65
            rr(draw, (xx, yy, xx+260, yy+105), 14, '#F4E4C2', c, 4)
            label(draw, (xx+130, yy+53), t, 21, c, True, True, 'mm')
        label(draw, (cx, y1+95), 'AI SERVICE PACK', 28, INK, True, True, 'mm')
    elif scene == 'gateway':
        rr(draw, (cx-120, cy-250, cx+120, cy+260), 18, '#6A5848', GOLD, 7)
        label(draw, (cx, cy), 'API\nGATEWAY', 25, PAPER_LIGHT, True, True, 'mm')
        for i, (t, c) in enumerate([('OCR', BLUE), ('VISION', GREEN), ('ASR', ORANGE), ('TTS', RED)]):
            yy = y1+80+i*140
            workshop(draw, x2-270, yy, t, c, .85)
            arrow(draw, (cx+130, cy), (x2-285, yy+60), c, 5)
        for i in range(3):
            arrow(draw, (x1+20, cy-130+i*130), (cx-130, cy-100+i*100), BLUE, 5)
    elif scene == 'token':
        for i, (role, c) in enumerate([('DEV', BLUE), ('PM', ORANGE), ('APP', GREEN)]):
            person(draw, x1+120+i*220, cy+90, .65, c, True, role)
            rr(draw, (x1+45+i*220, cy-220, x1+195+i*220, cy-115), 14, '#FFF4D7', c, 4)
            label(draw, (x1+120+i*220, cy-168), 'TOKEN', 18, c, True, True, 'mm')
        rr(draw, (x2-250, cy-180, x2-50, cy+180), 24, '#5B5043', GOLD, 6)
        label(draw, (x2-150, cy), 'ENTRY', 22, PAPER_LIGHT, True, True, 'mm')
    elif scene == 'usage':
        rr(draw, (cx-330, cy-230, cx+330, cy+230), 28, '#F6E8C8', BLUE, 5)
        for i, (t, c, v) in enumerate([('WHO', BLUE, .72), ('MODE', GREEN, .56), ('COUNT', ORANGE, .84), ('RESULT', RED, .38)]):
            yy = cy-155+i*95
            label(draw, (cx-270, yy), t, 20, c, True, True, 'lm')
            draw.rounded_rectangle((cx-120, yy-16, cx+250, yy+16), radius=12, fill='#D8C9AA')
            draw.rounded_rectangle((cx-120, yy-16, cx-120+370*v, yy+16), radius=12, fill=c)
    elif scene == 'queue':
        for i in range(6):
            xx = x1+65+i*120
            rr(draw, (xx, cy-95+(i%2)*30, xx+95, cy+35+(i%2)*30), 14, '#F3E4C3', BLUE if i < 4 else ORANGE, 4)
            label(draw, (xx+48, cy-30+(i%2)*30), f'T-{104+i}', 17, INK, True, True, 'mm')
        arrow(draw, (x1+30, cy+125), (x2-250, cy+125), GREEN, 8)
        person(draw, x2-135, cy-40, .7, GREEN, False, 'WORK')
    elif scene == 'artifact':
        person(draw, x1+120, cy-30, .7, GREEN, False, 'WORK')
        arrow(draw, (x1+220, cy), (cx-50, cy), BLUE, 8)
        for i, (t, c) in enumerate([('WAV', ORANGE), ('PNG', BLUE), ('JSON', GREEN), ('MODEL', RED)]):
            rr(draw, (cx+i%2*210, cy-160+i//2*190, cx+170+i%2*210, cy-40+i//2*190), 15, '#F8EBCD', c, 4)
            label(draw, (cx+85+i%2*210, cy-100+i//2*190), t, 20, c, True, True, 'mm')
        label(draw, (cx+190, cy+280), 'ARTIFACT SHELF / RETENTION', 20, MUTED, True, True, 'mm')
    elif scene == 'hub_core':
        draw.ellipse((cx-145, cy-145, cx+145, cy+145), fill='#D9E7E5', outline=BLUE, width=9)
        label(draw, (cx, cy), '3waAIHub', 30, INK, True, True, 'mm')
        parts = [('SERVICE', BLUE), ('TASK', ORANGE), ('ACCOUNT', GREEN), ('TOKEN', GOLD), ('USAGE', RED), ('HEALTH', METAL)]
        for i, (t, c) in enumerate(parts):
            a = -90+i*60
            px = cx+math.cos(math.radians(a))*340
            py = cy+math.sin(math.radians(a))*260
            line(draw, [(cx, cy), (px, py)], c, 5)
            rr(draw, (px-80, py-38, px+80, py+38), 16, '#FFF6DF', c, 4)
            label(draw, (px, py), t, 17, c, True, True, 'mm')
    elif scene == 'grid':
        hubs = [('VISION', BLUE), ('OCR', GREEN), ('LLM', RED), ('ASR', ORANGE), ('TTS', GOLD), ('YOLO', BLUE), ('SAM', GREEN), ('BIREf', ORANGE), ('DOC', METAL)]
        draw.ellipse((cx-120, cy-120, cx+120, cy+120), fill='#D9E6E4', outline=BLUE, width=9)
        label(draw, (cx, cy), 'AIHUB\nGRID', 25, INK, True, True, 'mm')
        for i, (t, c) in enumerate(hubs):
            a = -90+i*40
            px = cx+math.cos(math.radians(a))*390
            py = cy+math.sin(math.radians(a))*300
            line(draw, [(cx, cy), (px, py)], BLUE_LIGHT, 4)
            workshop(draw, px-70, py-55, t, c, .72)
    elif scene == 'long_jobs':
        for i, (t, c) in enumerate([('TRAIN', RED), ('PARSE', BLUE), ('VOICE', ORANGE), ('CONVERT', GREEN)]):
            xx = x1+45+(i%2)*390
            yy = y1+50+(i//2)*290
            workshop(draw, xx, yy, t, c, 1.05)
            draw.arc((xx+230, yy+20, xx+330, yy+120), 20, 330, fill=c, width=7)
    elif scene == 'browser_wait':
        rr(draw, (x1+50, cy-220, x1+410, cy+150), 22, '#EEE6D2', BROWN, 5)
        line(draw, [(x1+50, cy-150), (x1+410, cy-150)], BROWN, 4)
        label(draw, (x1+230, cy-30), 'WAITING…', 30, RED, True, True, 'mm')
        arrow(draw, (x1+430, cy), (x2-350, cy), BLUE, 8)
        workshop(draw, x2-300, cy-100, 'BACKGROUND\nWORK', GREEN, 1.1)
    elif scene == 'queue_flow':
        stages = [('SUBMIT', BLUE), ('QUEUE', GOLD), ('WORKER', ORANGE), ('RESULT', GREEN)]
        for i, (t, c) in enumerate(stages):
            xx = x1+30+i*200
            rr(draw, (xx, cy-90, xx+155, cy+90), 18, '#F7EACD', c, 5)
            label(draw, (xx+78, cy), t, 19, c, True, True, 'mm')
            if i < 3:
                arrow(draw, (xx+165, cy), (xx+195, cy), c, 6)
        label(draw, (cx, cy+180), 'TASK ID：T-2048', 22, MUTED, True, True, 'mm')
    elif scene == 'taskflow':
        stages = [('QUEUE', BLUE), ('LEASE', GOLD), ('WORKER', ORANGE), ('HEARTBEAT', GREEN), ('RESULT', BLUE)]
        for i, (t, c) in enumerate(stages):
            xx = x1+10+i*160
            draw.ellipse((xx, cy-65, xx+130, cy+65), fill='#F6E8C8', outline=c, width=5)
            label(draw, (xx+65, cy), t, 15, c, True, True, 'mm')
            if i < 4:
                arrow(draw, (xx+135, cy), (xx+155, cy), c, 5)
        for i, t in enumerate(['TIMEOUT', 'CANCEL', 'RECOVERY']):
            label(draw, (x1+180+i*220, cy+160), t, 18, RED, True, True, 'mm')
    elif scene == 'yolo':
        workshop(draw, cx-100, cy-80, 'YOLO', BLUE, 1.15)
        for i, (t, c) in enumerate([('PREDICT', GREEN), ('TRAIN', ORANGE), ('ONNX EXPORT', RED)]):
            px = x1+80+i*260
            py = cy+240
            arrow(draw, (cx, cy+80), (px+95, py-20), c, 6)
            rr(draw, (px, py, px+190, py+90), 15, '#F8EBCD', c, 4)
            label(draw, (px+95, py+45), t, 17, c, True, True, 'mm')
    elif scene == 'workers':
        for i, (t, c) in enumerate([('GPU', ORANGE), ('CPU', GREEN), ('DOC', BLUE), ('MEDIA', RED)]):
            person(draw, x1+100+i*190, cy+20+(i%2)*30, .62, c, False, t)
        rr(draw, (cx-160, y1+20, cx+160, y1+160), 22, '#D8E6E4', BLUE, 5)
        label(draw, (cx, y1+90), 'TASK CONTROL', 22, BLUE, True, True, 'mm')
    elif scene == 'overheat':
        server(draw, cx-120, cy-150, 240, 320, ORANGE, 'GPU / VRAM')
        for i in range(5):
            draw.arc((cx-20+i*28, cy-300-i*18, cx+150+i*35, cy-80+i*15), 180, 320, fill=(215,101,50,190-i*22), width=10)
        label(draw, (cx+280, cy), '95%', 70, RED, True, True, 'mm')
    elif scene == 'bigger':
        server(draw, cx-70, cy-250, 300, 410, BLUE, 'BIGGER GPU')
        person(draw, x1+140, cy+40, .75, BLUE, True, '?')
        storyteller(draw, x2-150, cy, .7)
        bubble(draw, (x1+260, y1+30, x2-290, y1+230), '換一台更大的？', BLUE, 36)
    elif scene == 'limits':
        server(draw, x1+100, cy-160, 230, 330, BLUE, 'NEW GPU')
        for i, (t, c) in enumerate([('LLM', RED), ('VOICE', ORANGE), ('VISION', GREEN)]):
            model_box(draw, x1+360+i*125, cy-180+i*80, 210, 125, t, c)
        server(draw, x2-210, cy+80, 150, 210, METAL, 'IDLE')
    elif scene == 'three_nodes':
        nodes = [('GTX 1080', '8 GB', BLUE), ('RTX 5090', '32 GB', ORANGE), ('RTX 5060 Ti', '16 GB', GREEN)]
        for i, (name, cap, c) in enumerate(nodes):
            xx = x1+20+i*260
            server(draw, xx, cy-150+(i%2)*45, 180, 255, c, name)
            label(draw, (xx+90, cy+160+(i%2)*45), cap, 28, c, True, True, 'mm')
        label(draw, (cx, y2-70), 'ROUTE BY CAPABILITY — NOT VRAM POOLING', 18, MUTED, True, True, 'mm')
    elif scene == 'cluster':
        rr(draw, (cx-100, cy-90, cx+100, cy+90), 22, '#D8E7E5', BLUE, 6)
        label(draw, (cx, cy), 'ROUTER', 22, BLUE, True, True, 'mm')
        for i, (t, c) in enumerate([('VISION', BLUE), ('VOICE', ORANGE), ('TRAIN', RED), ('CPU TOOLS', GREEN)]):
            a = -120+i*80
            px = cx+math.cos(math.radians(a))*350
            py = cy+math.sin(math.radians(a))*250
            server(draw, px-65, py-80, 130, 170, c, t)
            arrow(draw, (cx, cy), (px, py), c, 5)
    elif scene == 'machine_room':
        rr(draw, (x1+20, cy-170, x1+190, cy+170), 22, '#665445', GOLD, 6)
        label(draw, (x1+105, cy), 'ONE\nENTRY', 21, PAPER_LIGHT, True, True, 'mm')
        for row in range(2):
            for col in range(3):
                xx = x1+340+col*190
                yy = y1+40+row*300
                server(draw, xx, yy, 130, 205, [BLUE, ORANGE, GREEN][col], f'N{row*3+col+1}')
                arrow(draw, (x1+200, cy), (xx, yy+95), BLUE_LIGHT, 4)
    elif scene == 'islands':
        for i, (t, c) in enumerate([('ENGINEER', BLUE), ('PROJECT', ORANGE), ('DEPT A', GREEN), ('DEPT B', RED)]):
            px = x1+100+(i%2)*390
            py = y1+100+(i//2)*270
            draw.ellipse((px-90, py-45, px+180, py+100), fill='#CFB888', outline=c, width=5)
            model_box(draw, px, py-65, 120, 85, t, c)
    elif scene == 'exchange':
        rr(draw, (cx-150, cy-120, cx+150, cy+120), 24, '#D8E7E5', BLUE, 6)
        label(draw, (cx, cy), 'AI 能力\n交換所', 26, INK, True, anchor='mm')
        for i, (t, c) in enumerate([('DEPT A', BLUE), ('DEPT B', ORANGE), ('PROJECT', GREEN), ('USER', RED)]):
            a = -135+i*90
            px = cx+math.cos(math.radians(a))*360
            py = cy+math.sin(math.radians(a))*260
            workshop(draw, px-70, py-55, t, c, .72)
            arrow(draw, (px, py), (cx, cy), c, 5)
    elif scene == 'architecture':
        # Full platform map, still drawn as a pop-up technical city.
        top = [('工程師', BLUE), ('系統', GREEN), ('Agent', ORANGE)]
        for i, (t, c) in enumerate(top):
            rr(draw, (x1+20+i*185, y1+10, x1+175+i*185, y1+80), 14, '#FFF5DD', c, 4)
            label(draw, (x1+98+i*185, y1+45), t, 18, c, True, anchor='mm')
        rr(draw, (x1+185, y1+125, x1+590, y1+225), 18, '#E4D4AF', BLUE, 5)
        label(draw, (x1+388, y1+175), '帳號／Token → API Gateway', 20, INK, True, False, 'mm')
        layers = [
            ('LOCAL SERVICE', ['LLM', 'OCR', 'ASR', 'TTS'], BLUE),
            ('LOCAL EXECUTION', ['QUEUE', 'WORKER', 'LEASE', 'ARTIFACT'], ORANGE),
            ('CLUSTER', ['ROUTER', 'STATION', 'HEALTH', 'VRAM'], GREEN),
        ]
        for i, (t, parts, c) in enumerate(layers):
            xx = x1+20+i*270
            yy = y1+290
            rr(draw, (xx, yy, xx+235, yy+250), 20, '#F7E9CB', c, 5)
            label(draw, (xx+118, yy+38), t, 16, c, True, True, 'mm')
            for j, p in enumerate(parts):
                label(draw, (xx+118, yy+85+j*38), p, 16, INK, False, True, 'mm')
            arrow(draw, (x1+388, y1+230), (xx+118, yy-5), c, 4)
        label(draw, (cx, y2-55), 'USAGE · LOGS · CALLBACK · RETENTION · GOVERNANCE', 18, MUTED, True, True, 'mm')
    elif scene == 'city':
        # Dense final city/grid: departments and capabilities form one governed network.
        for i in range(16):
            col, row = i%8, i//8
            xx = x1+20+col*105
            yy = y1+130+row*270+(col%2)*28
            h = 85+(i*37)%110
            draw.rectangle((xx, yy-h, xx+65, yy), fill=['#79634E', '#687272', '#8A6A4C'][i%3], outline=INK, width=3)
            for k in range(3):
                draw.rectangle((xx+12, yy-h+18+k*24, xx+24, yy-h+30+k*24), fill=BLUE_LIGHT)
        rr(draw, (cx-125, cy-100, cx+125, cy+100), 24, '#D8E7E5', BLUE, 7)
        label(draw, (cx, cy), '3waAIHub', 27, INK, True, True, 'mm')
        for i in range(12):
            a = i*30
            px = cx+math.cos(math.radians(a))*420
            py = cy+math.sin(math.radians(a))*260
            line(draw, [(cx, cy), (px, py)], BLUE if i%2==0 else GREEN, 4)
            draw.ellipse((px-10, py-10, px+10, py+10), fill=ORANGE if i%3==0 else BLUE)
    elif scene == 'qa':
        storyteller(draw, x1+150, cy+20, .78)
        rr(draw, (cx-220, cy-170, cx+250, cy+180), 24, '#795438', GOLD, 7)
        line(draw, [(cx-205, cy), (cx+235, cy)], '#D7BB85', 5)
        for i in range(3):
            person(draw, x2-330+i*115, cy+120+(i%2)*15, .58, [BLUE, ORANGE, GREEN][i], True)
        line(draw, [(x2-215, cy+50), (x2-215, cy-150)], BLUE, 7)
        draw.ellipse((x2-230, cy-175, x2-200, cy-145), fill=BLUE)


def render_slide(item: dict) -> Image.Image:
    img = paper_canvas(item['page'] * 7919)
    shadow = Image.new('RGBA', (W, H), (0, 0, 0, 0))
    art = Image.new('RGBA', (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(art, 'RGBA')
    sd = ImageDraw.Draw(shadow, 'RGBA')
    page = item['page']
    chapter = item['chapter']
    # Alternating book-page focus preserves the feeling of a continuous flip.
    left_text = page % 2 == 1
    if item['scene'] in {'cover', 'answer', 'hub_reveal', 'grid', 'architecture', 'city', 'qa'}:
        left_text = True
    text_box = (92, 165, 855, 910) if left_text else (1060, 165, 1828, 910)
    art_box = (990, 165, 1815, 900) if left_text else (105, 165, 930, 900)
    # Chapter tab and folio.
    tab_x = 92 if left_text else 1060
    rr(d, (tab_x, 76, tab_x+300, 126), 16, '#E3CFA4', GOLD, 2, sd)
    label(d, (tab_x+18, 101), chapter, 18, BROWN, True, False, 'lm')
    label(d, (1825, 1007), f'{page:02d} / 48', 18, MUTED, True, True, 'ra')
    label(d, (95, 1007), '3waAIHub · GIDAY 08.07', 16, MUTED, False, True)

    if item['scene'] == 'cover':
        draw_scene(d, item['scene'], (900, 120, 1800, 930), page)
        d.multiline_text((105, 205), item['main'], font=font(SERIF_BOLD, 72), fill=INK, spacing=18)
        label(d, (110, 460), '3waAIHub：從單機模型、服務盒，\n到跨主機 AI 數位治理平台', 31, MUTED, True)
        label(d, (110, 820), 'GIS 地理資訊中心 GIDAY 分享會', 24, BLUE, True)
        label(d, (110, 865), 'AI 算了吧...', 22, GREEN, True)
    elif item['scene'] in {'answer', 'grid', 'architecture', 'city', 'qa'}:
        label(d, (105, 165), item['title'], 50, INK, True)
        if item['scene'] == 'architecture':
            d.multiline_text((105, 238), item['main'], font=font(SERIF_BOLD, 42), fill=INK, spacing=10)
            draw_scene(d, item['scene'], (820, 175, 1800, 900), page)
        elif item['scene'] == 'city':
            d.multiline_text((105, 250), item['main'], font=font(SERIF_BOLD, 54), fill=INK, spacing=12)
            draw_scene(d, item['scene'], (820, 180, 1810, 910), page)
            label(d, (108, 610), '從工具，走向服務。', 27, BLUE, True)
            label(d, (108, 655), '從服務，走向平台。', 27, ORANGE, True)
            label(d, (108, 700), '從平台，走向治理。', 27, GREEN, True)
            label(d, (108, 835), '不是每個人都要蓋一座發電廠，\n但每個人都應該能使用電。', 25, BROWN, True)
        elif item['scene'] == 'qa':
            d.multiline_text((105, 250), item['main'], font=font(SERIF_BOLD, 70), fill=INK, spacing=14)
            draw_scene(d, item['scene'], (870, 160, 1810, 910), page)
            label(d, (108, 810), '3waAIHub · AI 數位治理服務平台', 24, BLUE, True)
        else:
            d.multiline_text((105, 245), item['main'], font=font(SERIF_BOLD, 58), fill=INK, spacing=13)
            draw_scene(d, item['scene'], (890, 165, 1805, 920), page)
    else:
        tx1, ty1, tx2, ty2 = text_box
        label(d, (tx1, ty1), item['title'], 50, INK, True)
        d.multiline_text((tx1, ty1+110), item['main'], font=font(SERIF_BOLD, 57), fill=INK, spacing=16)
        draw_scene(d, item['scene'], art_box, page)

    # Fine blue route crossing the spine provides continuity from page to page.
    route_y = 944 + int(math.sin(page*0.72)*13)
    d.line((105, route_y, 930, route_y-3, 990, route_y+3, 1815, route_y), fill=(29,131,181,95), width=3)
    d.ellipse((950, route_y-10, 970, route_y+10), fill=(215,101,50,180))
    img = Image.alpha_composite(img, shadow.filter(ImageFilter.GaussianBlur(12)))
    img = Image.alpha_composite(img, art)
    return img.convert('RGB')


def build() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = '從一個工程師的 AI 冒險，到一座 AI 能力電網'
    prs.core_properties.subject = '3waAIHub：從單機模型、服務盒，到跨主機 AI 數位治理平台'
    prs.core_properties.author = 'AI 算了吧...'
    prs.core_properties.keywords = '3waAIHub, GIDAY, AI治理, Cluster, GPU, Local Execution'
    blank = prs.slide_layouts[6]

    assert len(SLIDES) == 48
    for item in SLIDES:
        image = render_slide(item)
        path = ASSETS / f'slide-{item["page"]:02d}.png'
        image.save(path, optimize=True)
        if item['page'] == 1:
            image.save(MASTER, optimize=True)
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
        notes = slide.notes_slide.notes_text_frame
        note_text = (
            f'【頁碼】{item["page"]:02d}\n'
            f'【頁面標題】{item["title"]}\n'
            f'【畫面主要文字】{item["main"]}\n\n'
            f'【畫面構圖】{item["composition"]}\n\n'
            f'【講者旁白】{item["narration"]}\n\n'
            f'【下一頁銜接】{item["transition"]}'
        )
        if item['source']:
            note_text += f'\n\n【資料來源】{item["source"]}'
        notes.text = note_text

    assert len(prs.slides) == 48
    assert all(slide.notes_slide.notes_text_frame.text.strip() for slide in prs.slides)
    prs.save(OUT)
    print(OUT)


if __name__ == '__main__':
    build()
